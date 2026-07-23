from __future__ import annotations

from copy import deepcopy
import os
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
import zipfile
from lxml import etree
from PIL import Image as PILImage
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE_DASH_STYLE
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

from . import design as d

# Image formats that python-pptx/Pillow can hand to add_picture() without
# conversion. EMF is embedded as image/x-emf via _install_emf_image_support()
# (Pillow otherwise reports it as WMF, which PowerPoint renders as a blank box).
_NATIVE_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".wmf", ".emf"}


ASSETS = None

# Namespace URIs used for Office Math compatibility wrappers.
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _install_emf_image_support() -> None:
    """Teach python-pptx to embed EMF metafiles as ``image/x-emf``.

    python-pptx derives an image's part extension and content type from the
    format Pillow reports, and Pillow reports *both* WMF and EMF metafiles as
    ``"WMF"``. So an EMF is embedded as an ``image/x-wmf`` part named
    ``imageN.wmf`` -- PowerPoint then tries to parse the EMF bytes as WMF and
    renders nothing (a blank/black box). EMF files carry the ``" EMF"``
    signature at byte offset 40, so detect that and return the correct ``emf``
    extension (``image_content_types`` already maps ``emf`` -> ``image/x-emf``).
    """
    try:
        from pptx.parts.image import Image as _Image
    except Exception:
        return
    # Canonical extensions python-pptx assigns per Pillow format (mirrors its own
    # table); EMF is handled ahead of this by signature since Pillow calls it WMF.
    _ext_map = {"BMP": "bmp", "GIF": "gif", "JPEG": "jpg",
                "PNG": "png", "TIFF": "tiff", "WMF": "wmf"}

    def _ext(self):
        blob = self._blob
        if len(blob) >= 44 and blob[40:44] == b" EMF":
            return "emf"
        fmt = self._format
        if fmt not in _ext_map:
            raise ValueError("unsupported image format, got %r" % (fmt,))
        return _ext_map[fmt]

    # content_type reads self.ext (image_content_types[self.ext]), so patching
    # ext alone is enough to also correct the content type.
    _Image.ext = property(_ext)


_install_emf_image_support()


def set_assets_dir(path: Path | str) -> None:
    """Set the directory where image/movie assets are stored.

    Args:
        path: Directory path containing image and video assets. May be a
            ``pathlib.Path`` or a string.

    Returns:
        None.
    """
    global ASSETS
    ASSETS = Path(path)


def rgb(hex_str: str) -> RGBColor:
    """Convert a hex color string to an ``RGBColor``.

    Args:
        hex_str: RGB hex string such as ``'FF0000'`` or ``'#FF0000'``.

    Returns:
        A python-pptx ``RGBColor`` instance.
    """
    return RGBColor.from_string(hex_str.lstrip("#"))


# ---------------------------------------------------------------------------
# Theme / color helpers
# ---------------------------------------------------------------------------

THEME_COLOR_MAP = {
    "theme_accent1": MSO_THEME_COLOR.ACCENT_1,
    "theme_accent2": MSO_THEME_COLOR.ACCENT_2,
    "theme_accent3": MSO_THEME_COLOR.ACCENT_3,
    "theme_accent4": MSO_THEME_COLOR.ACCENT_4,
    "theme_accent5": MSO_THEME_COLOR.ACCENT_5,
    "theme_accent6": MSO_THEME_COLOR.ACCENT_6,
    "theme_bg1": MSO_THEME_COLOR.BACKGROUND_1,
    "theme_bg2": MSO_THEME_COLOR.BACKGROUND_2,
    "theme_tx1": MSO_THEME_COLOR.TEXT_1,
    "theme_tx2": MSO_THEME_COLOR.TEXT_2,
    "theme_lt1": MSO_THEME_COLOR.LIGHT_1,
    "theme_lt2": MSO_THEME_COLOR.LIGHT_2,
    "theme_dk1": MSO_THEME_COLOR.DARK_1,
    "theme_dk2": MSO_THEME_COLOR.DARK_2,
    "theme_hlink": MSO_THEME_COLOR.HYPERLINK,
    "theme_folHlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
}


# Common DrawingML preset colors mapped to RGB hex.
PRESET_COLOR_MAP = {
    "black": "000000",
    "white": "FFFFFF",
    "red": "FF0000",
    "green": "00FF00",
    "blue": "0000FF",
    "yellow": "FFFF00",
    "cyan": "00FFFF",
    "magenta": "FF00FF",
    "darkRed": "8B0000",
    "darkGreen": "006400",
    "darkBlue": "00008B",
    "darkYellow": "808000",
    "gray": "808080",
    "lightGray": "D3D3D3",
    "darkGray": "A9A9A9",
    "orange": "FFA500",
    "purple": "800080",
    "brown": "A52A2A",
    "pink": "FFC0CB",
    "lime": "00FF00",
    "teal": "008080",
    "navy": "000080",
    "maroon": "800000",
    "olive": "808000",
    "silver": "C0C0C0",
}


def _is_theme(color):
    if isinstance(color, dict):
        color = color.get("color", "")
    return isinstance(color, str) and (
        color.startswith("theme_") or color.startswith("scheme:")
    )


def _theme_color_value(color: str):
    """Return the MSO_THEME_COLOR enum value for a theme color string."""
    return THEME_COLOR_MAP.get(color)


def _preset_color_value(color: str):
    """Return the RGB hex value for a preset color string like 'prst:black'."""
    if isinstance(color, str) and color.startswith("prst:"):
        return PRESET_COLOR_MAP.get(color.split(":", 1)[1])
    return None


# OOXML luminance/saturation/tint/shade transforms, in a stable emit order. A
# light color is typically a darker base plus a lightening transform, so these
# must be re-emitted or the color reverts to its darker base. See the matching
# ``_COLOR_MOD_TAGS`` in ``scripts/helpers/slide_xml.py``.
_COLOR_MOD_TAGS = ("lumMod", "lumOff", "satMod", "satOff", "hueMod", "hueOff", "shade", "tint")


def _color_mods(color):
    """Return the {tag: val} transform map carried on a color dict, or {}."""
    if not isinstance(color, dict):
        return {}
    return {tag: color[tag] for tag in _COLOR_MOD_TAGS if color.get(tag) is not None}


def _apply_color_transforms(color_el, mods):
    """Append luminance/saturation/tint/shade transform children to a color element.

    ``mods`` maps an OOXML transform tag to its raw ``val`` string, e.g.
    ``{"lumMod": "40000", "lumOff": "60000"}``.
    """
    for tag in _COLOR_MOD_TAGS:
        val = mods.get(tag)
        if val is not None:
            etree.SubElement(color_el, qn(f"a:{tag}")).set("val", str(val))


def _apply_theme_color(color_format, color):
    """Apply a theme color (e.g. theme_accent1), preset color, RGB hex, or a color
    dict (with alpha and/or luminance/tint transforms) to a ColorFormat object."""
    alpha = None
    mods = _color_mods(color)
    if isinstance(color, dict):
        alpha = color.get("alpha")
        color = color.get("color")
    if color is None:
        return
    tc = _theme_color_value(color)
    if tc is not None:
        color_format.theme_color = tc
    elif _preset_color_value(color):
        color_format.rgb = rgb(_preset_color_value(color))
    else:
        color_format.rgb = rgb(color)
    if alpha is not None or mods:
        # ColorFormat has no direct alpha/transform API; reach into the XML element.
        color_el = color_format._color._xClr
        _apply_color_transforms(color_el, mods)
        if alpha is not None:
            etree.SubElement(color_el, qn("a:alpha")).set("val", str(int(round(alpha * 100000))))


def _apply_color_element(parent, color, shade=None, tint=None):
    """Append an <a:schemeClr> or <a:srgbClr> child to ``parent``.

    Luminance/tint/shade transforms are taken from the color dict when present.
    The explicit ``shade``/``tint`` params remain for backward compatibility with
    previously generated code and, when given, take precedence over the dict.
    """
    alpha = None
    mods = _color_mods(color)
    if isinstance(color, dict):
        alpha = color.get("alpha")
        color = color.get("color")
    if shade is not None:
        mods["shade"] = shade
    if tint is not None:
        mods["tint"] = tint
    if _is_theme(color):
        c = etree.SubElement(parent, qn("a:schemeClr"))
        if color.startswith("theme_"):
            c.set("val", color.replace("theme_", ""))
        else:
            c.set("val", color.split(":", 1)[1])
    else:
        c = etree.SubElement(parent, qn("a:srgbClr"))
        pc = _preset_color_value(color)
        c.set("val", pc if pc is not None else color.lstrip("#").upper())
    _apply_color_transforms(c, mods)
    if alpha is not None:
        etree.SubElement(c, qn("a:alpha")).set("val", str(int(round(alpha * 100000))))


# ---------------------------------------------------------------------------
# Shape kind mapping
# ---------------------------------------------------------------------------

_SHAPE_KIND_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "rightTriangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "chevron": MSO_SHAPE.CHEVRON,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "diamond": MSO_SHAPE.DIAMOND,
    "cross": MSO_SHAPE.CROSS,
    "heart": MSO_SHAPE.HEART,
    "star": MSO_SHAPE.STAR_5_POINT,
    "star5": MSO_SHAPE.STAR_5_POINT,
    "star6": MSO_SHAPE.STAR_6_POINT,
    "star8": MSO_SHAPE.STAR_8_POINT,
    "star10": MSO_SHAPE.STAR_10_POINT,
    "star12": MSO_SHAPE.STAR_12_POINT,
    "star16": MSO_SHAPE.STAR_16_POINT,
    "star24": MSO_SHAPE.STAR_24_POINT,
    "star32": MSO_SHAPE.STAR_32_POINT,
    "star4": MSO_SHAPE.STAR_4_POINT,
    "star7": MSO_SHAPE.STAR_7_POINT,
    "rightArrow": MSO_SHAPE.RIGHT_ARROW,
    "leftArrow": MSO_SHAPE.LEFT_ARROW,
    "upArrow": MSO_SHAPE.UP_ARROW,
    "downArrow": MSO_SHAPE.DOWN_ARROW,
    "leftRightArrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
    "upDownArrow": MSO_SHAPE.UP_DOWN_ARROW,
    "bentArrow": MSO_SHAPE.BENT_ARROW,
    "curvedRightArrow": MSO_SHAPE.CURVED_RIGHT_ARROW,
    "curvedLeftArrow": MSO_SHAPE.CURVED_LEFT_ARROW,
    "line": MSO_SHAPE.LINE_INVERSE,
    "rightBrace": MSO_SHAPE.RIGHT_BRACE,
    "leftBrace": MSO_SHAPE.LEFT_BRACE,
    "sun": MSO_SHAPE.SUN,
    "cloud": MSO_SHAPE.CLOUD,
    "smileyFace": MSO_SHAPE.SMILEY_FACE,
    "noSymbol": MSO_SHAPE.NO_SYMBOL,
    "can": MSO_SHAPE.CAN,
    "cube": MSO_SHAPE.CUBE,
    "bevel": MSO_SHAPE.BEVEL,
    "foldedCorner": MSO_SHAPE.FOLDED_CORNER,
    "frame": MSO_SHAPE.FRAME,
    "plaque": MSO_SHAPE.PLAQUE,
    "donut": MSO_SHAPE.DONUT,
    "arc": MSO_SHAPE.ARC,
    "blockArc": MSO_SHAPE.BLOCK_ARC,
    "chord": MSO_SHAPE.CHORD,
    "pie": MSO_SHAPE.PIE,
    "teardrop": MSO_SHAPE.TEAR,
    "wave": MSO_SHAPE.WAVE,
    "doubleWave": MSO_SHAPE.DOUBLE_WAVE,
}


def _shape_kind(kind: str):
    if isinstance(kind, MSO_SHAPE):
        return kind
    if kind in _SHAPE_KIND_MAP:
        return _SHAPE_KIND_MAP[kind]
    # Fall back to resolving the raw OOXML preset name (e.g. 'corner',
    # 'homePlate', 'chord', ...). The map above only lists common presets and
    # a few intentional remaps (e.g. 'line'); anything else that is a valid
    # preset geometry should round-trip as itself rather than silently
    # degrading to a rectangle.
    if kind:
        try:
            return MSO_SHAPE.from_xml(kind)
        except Exception:
            pass
    return MSO_SHAPE.RECTANGLE


# ---------------------------------------------------------------------------
# XML shape property helpers
# ---------------------------------------------------------------------------

def _suppress_default_shadow(shape):
    """Remove the default shape style shadow that python-pptx applies."""
    sp = shape._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)
    # NOTE: Do not write an empty <a:effectLst/> here. Some renderers (e.g.
    # PowerPoint) interpret an empty effect list placed before <a:ln> as
    # cancelling the line, so the shape border disappears.


def _apply_style(shape, style):
    """Write a <p:style> element from a parsed style dictionary.

    The style element must appear after <p:spPr> and before <p:txBody>, so it
    is inserted immediately after spPr rather than appended to the shape.
    """
    sp = shape._element
    old = sp.find(qn("p:style"))
    if old is not None:
        sp.remove(old)
    if not style:
        return
    style_el = etree.Element(qn("p:style"))
    for ref_key, ref in style.items():
        if not ref:
            continue
        ref_el = etree.SubElement(style_el, qn(f"a:{ref_key}"))
        ref_el.set("idx", str(ref.get("idx", "0")))
        color = ref.get("color")
        if color is not None:
            _apply_color_element(ref_el, color, ref.get("shade"), ref.get("tint"))

    spPr = sp.find(qn("p:spPr"))
    if spPr is not None:
        spPr.addnext(style_el)
    else:
        sp.append(style_el)


def _apply_adjustments(shape, adjustments):
    """Set adjustment values on a preset geometry."""
    if not adjustments:
        return
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn("a:avLst"))
    avLst.clear()
    for name, fmla in adjustments:
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", name)
        gd.set("fmla", fmla)


def _apply_rotation_and_flip(shape, rotation, flip_h, flip_v):
    if rotation:
        shape.rotation = rotation
    # python-pptx exposes no flip_horizontal/flip_vertical setter, so assigning
    # those attributes is a silent no-op that drops the flip. Write the xfrm
    # flip attributes directly (same approach as add_connector).
    if flip_h or flip_v:
        xfrm = shape._element.spPr.get_or_add_xfrm()
        if flip_h:
            xfrm.set("flipH", "1")
        if flip_v:
            xfrm.set("flipV", "1")


def _remove_autofit(tf):
    """Remove any autofit element so the saved XML matches PowerPoint's default."""
    bodyPr = tf._element.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for tag in (qn("a:noAutofit"), qn("a:spAutoFit"), qn("a:normAutofit")):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)


def _set_autofit(tf, autofit):
    """Write noAutofit/spAutoFit/normAutofit into the text body."""
    if autofit is None:
        return
    bodyPr = tf._element.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    _remove_autofit(tf)
    if autofit == "noAutofit":
        bodyPr.append(etree.Element(qn("a:noAutofit")))
    elif autofit == "spAutoFit":
        bodyPr.append(etree.Element(qn("a:spAutoFit")))
    elif autofit == "normAutofit":
        bodyPr.append(etree.Element(qn("a:normAutofit")))


def _set_text_margins(tf, margins):
    """Set text frame margins from a dict like {'lIns': 0.1, ...}."""
    if not margins:
        return
    bodyPr = tf._element.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    emu_per_inch = 914400
    for key in ("lIns", "rIns", "tIns", "bIns"):
        val = margins.get(key)
        if val is not None:
            try:
                bodyPr.set(key, str(int(float(val) * emu_per_inch)))
            except Exception:
                pass


def _set_wrap(tf, wrap):
    """Set word wrap and bodyPr/@wrap.

    ``wrap=None`` leaves the text frame's default wrap behaviour untouched.
    """
    if wrap is None:
        return
    tf.word_wrap = bool(wrap)
    bodyPr = tf._element.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("wrap", "square" if wrap else "none")


# ---------------------------------------------------------------------------
# Line / fill helpers
# ---------------------------------------------------------------------------

def _build_ln(color, width_pt=None, dash="solid", head=None, tail=None, cap="flat", cmpd="sng"):
    """Build a minimal <a:ln> element from scratch.

    Only writes attributes/children that differ from PowerPoint's defaults,
    keeping generated XML close to hand-authored files.
    """
    ln = etree.Element(qn("a:ln"))
    if width_pt is not None:
        ln.set("w", str(int(Pt(width_pt))))
    if cap and cap != "flat":
        ln.set("cap", cap)
    if cmpd and cmpd != "sng":
        ln.set("cmpd", cmpd)
    if color is False or color is None or color == "none":
        etree.SubElement(ln, qn("a:noFill"))
    else:
        sf = etree.SubElement(ln, qn("a:solidFill"))
        _apply_color_element(sf, color)
    if dash and dash != "solid":
        pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", dash)
    for attr, val in (("headEnd", head), ("tailEnd", tail)):
        if not val or val == "none":
            continue
        el = etree.SubElement(ln, qn(f"a:{attr}"))
        if isinstance(val, dict):
            end_type = val.get("type")
            if not end_type and (val.get("w") or val.get("len")):
                end_type = "arrow"
            if end_type:
                el.set("type", end_type)
            if val.get("w"):
                el.set("w", val["w"])
            if val.get("len"):
                el.set("len", val["len"])
        else:
            el.set("type", val)
    return ln


def _apply_line(shape, color=None, width_pt=None, dash="solid", head=None, tail=None, cap="flat", cmpd="sng"):
    """Replace the shape's <a:ln> with an explicit one.

    Passing ``color=None`` leaves the line untouched (no <a:ln> is added).
    Use ``color=False`` or ``color='none'`` to write an explicit no-fill line.
    """
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:ln")):
        spPr.remove(old)
    if color is None:
        return
    ln = _build_ln(color, width_pt, dash, head, tail, cap, cmpd)
    spPr.append(ln)


def _apply_fill(shape, fill):
    """Apply a fill to a shape or table cell.

    ``fill`` may be a color string, a theme color, ``None``/``False``/``'none'``,
    or a dict describing a gradient/pattern fill.
    """
    is_cell = not hasattr(shape, "_element")
    if is_cell:
        tc = shape._tc
    else:
        spPr = shape._element.spPr

    if fill is None:
        # Leave python-pptx's default. Style references are applied separately.
        return
    if fill is False or fill == "none":
        if is_cell:
            shape.fill.background()
        else:
            for tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:pattFill"), qn("a:blipFill")):
                for el in spPr.findall(tag):
                    spPr.remove(el)
            if spPr.find(qn("a:noFill")) is None:
                spPr.append(etree.Element(qn("a:noFill")))
        return
    if isinstance(fill, dict):
        if fill.get("type") == "gradient":
            shape.fill.gradient()
            angle = fill.get("angle", 0)
            if angle is not None:
                shape.fill.gradient_angle = angle
            stops = fill.get("stops", [])
            for i, stop in enumerate(stops):
                if i >= len(shape.fill.gradient_stops):
                    break
                gs = shape.fill.gradient_stops[i]
                gs.position = stop.get("pos", i / max(len(stops) - 1, 1))
                _apply_theme_color(gs.color, stop.get("color", "FFFFFF"))
            return
        if fill.get("type") == "pattern":
            shape.fill.patterned()
            tcPr = tc.get_or_add_tcPr() if is_cell else spPr
            pattFill = tcPr.find(qn("a:pattFill"))
            name = fill.get("name")
            if pattFill is not None and name:
                pattFill.set("prst", name)
            if fill.get("fg") is not None:
                _apply_theme_color(shape.fill.fore_color, fill["fg"])
            if fill.get("bg") is not None:
                _apply_theme_color(shape.fill.back_color, fill["bg"])
            return
    shape.fill.solid()
    _apply_theme_color(shape.fill.fore_color, fill)


# ---------------------------------------------------------------------------
# Effects helper
# ---------------------------------------------------------------------------

def _apply_effects(shape, effects):
    """Write an <a:effectLst> from a list of effect dicts."""
    if not effects:
        return
    spPr = shape._element.spPr
    old = spPr.find(qn("a:effectLst"))
    if old is not None:
        spPr.remove(old)
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    for eff in effects:
        kind = eff.get("type")
        if kind == "outerShdw":
            el = etree.SubElement(effectLst, qn("a:outerShdw"))
            for attr in ("blurRad", "dist", "dir", "algn", "rotWithShape"):
                if eff.get(attr) is not None:
                    el.set(attr, str(eff[attr]))
            _apply_color_element(el, eff.get("color", "000000"))
        elif kind == "innerShdw":
            el = etree.SubElement(effectLst, qn("a:innerShdw"))
            for attr in ("blurRad", "dist", "dir", "algn"):
                if eff.get(attr) is not None:
                    el.set(attr, str(eff[attr]))
            _apply_color_element(el, eff.get("color", "000000"))
        elif kind == "glow":
            el = etree.SubElement(effectLst, qn("a:glow"))
            el.set("rad", str(eff.get("rad", 40000)))
            _apply_color_element(el, eff.get("color", "FFFFFF"))
        elif kind == "reflection":
            el = etree.SubElement(effectLst, qn("a:reflection"))
            for attr in ("blurRad", "dist", "algn", "rotWithShape"):
                if eff.get(attr) is not None:
                    el.set(attr, str(eff[attr]))
        elif kind == "softEdge":
            el = etree.SubElement(effectLst, qn("a:softEdge"))
            el.set("rad", str(eff.get("rad", 50000)))


# ---------------------------------------------------------------------------
# Text / run helpers
# ---------------------------------------------------------------------------

def _apply_run(run, text, *, font=None, size=None, color=None, bold=None, italic=None,
               underline=None, strike=None, highlight=None, baseline=None,
               spacing=None, hyperlink=None, pitch_family=None, charset=None, effects=None):
    run.text = text
    if font is not None:
        run.font.name = font
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        _apply_theme_color(run.font.color, color)
    if bold is True:
        run.font.bold = True
    elif bold is False:
        rPr = run._r.get_or_add_rPr()
        rPr.attrib.pop("b", None)
    if italic is True:
        run.font.italic = True
    elif italic is False:
        rPr = run._r.get_or_add_rPr()
        rPr.attrib.pop("i", None)
    if underline is not None:
        rPr = run._r.get_or_add_rPr()
        if isinstance(underline, str):
            rPr.set("u", underline)
        elif underline:
            rPr.set("u", "sng")
        else:
            rPr.attrib.pop("u", None)
    if strike is not None:
        rPr = run._r.get_or_add_rPr()
        if isinstance(strike, str):
            rPr.set("strike", strike)
        elif strike:
            rPr.set("strike", "sngStrike")
        else:
            rPr.attrib.pop("strike", None)
    if baseline is not None:
        run._r.get_or_add_rPr().set("baseline", str(baseline))
    if spacing is not None:
        run._r.get_or_add_rPr().set("spc", str(int(spacing)))
    if highlight is not None:
        rPr = run._r.get_or_add_rPr()
        for old in rPr.findall(qn("a:highlight")):
            rPr.remove(old)
        # a:highlight must precede a:latin/underline/hlink in the CT_TextCharacterProperties
        # schema; insert at the correct position instead of appending after a:latin.
        hl = rPr.makeelement(qn("a:highlight"), {})
        rPr.insert_element_before(
            hl, "a:uLnTx", "a:uLn", "a:uFillTx", "a:uFill",
            "a:latin", "a:ea", "a:cs", "a:sym",
            "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")
        _apply_color_element(hl, highlight)
    if hyperlink is not None:
        run.hyperlink.address = hyperlink
    if pitch_family is not None or charset is not None:
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is not None:
                if pitch_family is not None:
                    el.set("pitchFamily", str(pitch_family))
                if charset is not None:
                    el.set("charset", str(charset))


def _set_paragraph_bullet(p, bullet):
    """Set paragraph bullet/no-bullet via raw pPr XML."""
    if bullet is None:
        return
    pPr = p._p.get_or_add_pPr()
    for tag in (qn("a:buClr"), qn("a:buSzPts"), qn("a:buSzPct"), qn("a:buFont"),
                qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum")):
        for el in pPr.findall(tag):
            pPr.remove(el)
    if bullet == "none":
        pPr.append(etree.Element(qn("a:buNone")))
    elif bullet == "char":
        bu = etree.Element(qn("a:buChar"))
        bu.set("char", "•")
        pPr.append(bu)
    elif bullet == "autoNum":
        bu = etree.Element(qn("a:buAutoNum"))
        bu.set("type", "arabicParenR")
        pPr.append(bu)
    elif isinstance(bullet, dict):
        kind = bullet.get("type")
        # Glyph formatting (color/size/font). _reorder_pPr moves these ahead of
        # the glyph element (buChar/buAutoNum) as the OOXML schema requires.
        color = bullet.get("color")
        if color is not None:
            _apply_color_element(etree.SubElement(pPr, qn("a:buClr")), color)
        if bullet.get("size_pts") is not None:
            etree.SubElement(pPr, qn("a:buSzPts")).set(
                "val", str(int(round(float(bullet["size_pts"]) * 100))))
        elif bullet.get("size_pct") is not None:
            etree.SubElement(pPr, qn("a:buSzPct")).set(
                "val", str(int(round(float(bullet["size_pct"]) * 1000))))
        if bullet.get("font") is not None:
            etree.SubElement(pPr, qn("a:buFont")).set("typeface", bullet["font"])
        if kind == "char":
            bu = etree.Element(qn("a:buChar"))
            bu.set("char", bullet.get("char", "•"))
            pPr.append(bu)
        elif kind == "autoNum":
            bu = etree.Element(qn("a:buAutoNum"))
            bu.set("type", bullet.get("style", "arabicParenR"))
            pPr.append(bu)
        elif kind == "none":
            pPr.append(etree.Element(qn("a:buNone")))


def _set_paragraph_spacing(p, line_spacing=None, space_before=None, space_after=None):
    pPr = p._p.get_or_add_pPr()
    if line_spacing is not None:
        for old in pPr.findall(qn("a:lnSpc")):
            pPr.remove(old)
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        if isinstance(line_spacing, str) and line_spacing.endswith("pts"):
            pts = line_spacing.replace("pts", "").strip()
            spcPts = etree.SubElement(lnSpc, qn("a:spcPts"))
            spcPts.set("val", str(int(float(pts) * 100)))
        else:
            spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
            spcPct.set("val", str(int(line_spacing)))
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)


def _resolve_enum(value, enum_cls):
    """Convert a string like 'MSO_ANCHOR.MIDDLE' to the actual enum member."""
    if value is None or isinstance(value, enum_cls):
        return value
    if isinstance(value, str):
        name = value.split(".")[-1]
        return getattr(enum_cls, name, value)
    return value


def _set_default_run_props(p, *, font=None, size=None, color=None, bold=None,
                           italic=None, underline=None, strike=None, highlight=None,
                           baseline=None, spacing=None, hyperlink=None, pitch_family=None, charset=None):
    """Write paragraph-level default run properties (<a:pPr><a:defRPr>).

    Default properties apply to runs that do not specify their own rPr,
    matching how PowerPoint stores single-style paragraphs.
    """
    pPr = p._p.get_or_add_pPr()
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))

    if size is not None:
        defRPr.set("sz", str(int(size * 100)))
    else:
        defRPr.attrib.pop("sz", None)
    if bold is True:
        defRPr.set("b", "1")
    else:
        defRPr.attrib.pop("b", None)
    if italic is True:
        defRPr.set("i", "1")
    else:
        defRPr.attrib.pop("i", None)
    if underline is not None:
        if isinstance(underline, str):
            defRPr.set("u", underline)
        elif underline:
            defRPr.set("u", "sng")
        else:
            defRPr.attrib.pop("u", None)
    if strike is not None:
        if isinstance(strike, str):
            defRPr.set("strike", strike)
        elif strike:
            defRPr.set("strike", "sngStrike")
        else:
            defRPr.attrib.pop("strike", None)
    if baseline:
        defRPr.set("baseline", str(baseline))
    else:
        defRPr.attrib.pop("baseline", None)
    if spacing is not None:
        defRPr.set("spc", str(int(spacing)))
    else:
        defRPr.attrib.pop("spc", None)
    if color is not None:
        for old in defRPr.findall(qn("a:solidFill")):
            defRPr.remove(old)
        sf = etree.SubElement(defRPr, qn("a:solidFill"))
        _apply_color_element(sf, color)
    if highlight is not None:
        for old in defRPr.findall(qn("a:highlight")):
            defRPr.remove(old)
        hl = etree.SubElement(defRPr, qn("a:highlight"))
        _apply_color_element(hl, highlight)
    if font is not None:
        latin = defRPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(defRPr, qn("a:latin"))
        latin.set("typeface", font)
        if pitch_family is not None:
            latin.set("pitchFamily", str(pitch_family))
        if charset is not None:
            latin.set("charset", str(charset))


_RUN_KEYS = ("font", "size", "color", "bold", "italic", "underline", "strike",
             "highlight", "baseline", "spacing", "hyperlink", "pitch_family", "charset")


_PPR_CHILD_ORDER = [
    "lnSpc", "spcBef", "spcAft",
    "buClr", "buSzPct", "buSzPts", "buFont", "buChar", "buAutoNum",
    "buNone", "buRTL",
    "tabLst", "defTabSz", "rtl", "eaLineBrk", "latinLineBrk", "hangingPunct",
    "defRPr", "extLst",
]


def _reorder_pPr(pPr):
    """Sort children of <a:pPr> into the order required by the OOXML schema."""
    order = {tag: idx for idx, tag in enumerate(_PPR_CHILD_ORDER)}

    def sort_key(child):
        tag = child.tag.split("}")[-1]
        # Unknown tags go to the end.
        return order.get(tag, len(_PPR_CHILD_ORDER))

    children = sorted(pPr, key=sort_key)
    pPr[:] = children


def _apply_text_frame(tf, paragraphs, *, font=None, size=None, color=None, bold=None,
                      italic=None, underline=None, strike=None, highlight=None, baseline=None,
                      hyperlink=None, align=None, anchor=None, vertical_anchor=None,
                      margins=None, wrap=None, autofit=None, line_spacing=None,
                      space_before=None, space_after=None, bullet=None,
                      pitch_family=None, charset=None):
    _set_wrap(tf, wrap)
    _remove_autofit(tf)
    va = _resolve_enum(vertical_anchor, MSO_ANCHOR) or _resolve_enum(anchor, MSO_ANCHOR)
    if va is not None:
        tf.vertical_anchor = va
    _set_text_margins(tf, margins)
    _set_autofit(tf, autofit)

    top_defaults = {
        "font": font, "size": size, "color": color, "bold": bold,
        "italic": italic, "underline": underline, "strike": strike,
        "highlight": highlight, "baseline": baseline,
        "pitch_family": pitch_family, "charset": charset,
    }

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(para, str):
            para = {"text": para}
        elif isinstance(para, tuple):
            text, level = para
            para = {"text": text, "level": level}
        text = para.get("text", "")
        runs = para.get("runs") if isinstance(para.get("runs"), list) else (
            text if isinstance(text, list) else None
        )

        lvl = para.get("level")
        if lvl is not None:
            p.level = lvl

        pa = _resolve_enum(para.get("align", align), PP_ALIGN)
        if pa is not None:
            p.alignment = pa
        else:
            pPr = p._p.find(qn("a:pPr"))
            if pPr is not None:
                pPr.attrib.pop("algn", None)

        pPr = p._p.get_or_add_pPr()
        marL = para.get("marL")
        indent = para.get("indent")
        if marL is not None:
            pPr.set("marL", str(marL))
        if indent is not None:
            pPr.set("indent", str(indent))

        para_bullet = para.get("bullet", bullet)
        _set_paragraph_bullet(p, para_bullet)

        para_line_spacing = para.get("line_spacing", line_spacing)
        para_space_before = para.get("space_before", space_before)
        para_space_after = para.get("space_after", space_after)
        _set_paragraph_spacing(p, para_line_spacing, para_space_before, para_space_after)

        # Paragraph-level default run formatting: top-level values overridden by
        # paragraph dict values. Runs below use only their own explicit props.
        defaults = dict(top_defaults)
        for key in _RUN_KEYS:
            if para.get(key) is not None:
                defaults[key] = para[key]
        # Boolean False has no effect on a defRPr (it means "not this attribute"),
        # so drop False values to avoid emitting empty default-run properties.
        defaults = {k: v for k, v in defaults.items() if v is not None and v is not False}
        if defaults:
            _set_default_run_props(p, **defaults)

        # Ensure paragraph properties children are in schema order (line spacing
        # and spacing before/after must precede bullet properties, and defRPr
        # must come after them).
        _reorder_pPr(p._p.get_or_add_pPr())

        last_run_el = None
        if runs:
            for j, rspec in enumerate(runs):
                if isinstance(rspec, str):
                    rspec = {"text": rspec}
                if rspec.get("math_xml"):
                    math_el = etree.fromstring(rspec["math_xml"])
                    if last_run_el is not None:
                        last_run_el.addnext(math_el)
                    else:
                        pPr = p._p.find(qn("a:pPr"))
                        if pPr is not None:
                            pPr.addnext(math_el)
                        else:
                            p._p.insert(0, math_el)
                    continue
                run = p.runs[0] if j == 0 and p.runs else p.add_run()
                run_args = {"text": rspec.get("text", "")}
                for key in _RUN_KEYS:
                    if rspec.get(key) is not None:
                        run_args[key] = rspec[key]
                _apply_run(run, **run_args)
                last_run_el = run._r
        else:
            if text != "":
                run = p.runs[0] if p.runs else p.add_run()
                _apply_run(run, text)
                last_run_el = run._r

        # Remove any empty placeholder runs created by python-pptx but not used.
        for r_el in list(p._p.findall(qn("a:r"))):
            t = r_el.find(qn("a:t"))
            if t is None or not t.text:
                p._p.remove(r_el)


# ---------------------------------------------------------------------------
def renumber_slide_ids(slide: "Slide") -> None:
    """Reassign unique sequential shape IDs across the entire slide tree.

    cNvPr elements inside ``mc:Fallback`` are skipped so the fallback shape keeps
    the same ID as its matching ``Choice`` shape, matching PowerPoint's own output.

    Args:
        slide: A python-pptx ``Slide`` whose shapes should be renumbered.

    Returns:
        None. IDs are updated in place on ``slide.element``.
    """
    ids = []
    for cNvPr in slide.element.iter(f"{{{P_NS}}}cNvPr"):
        # Skip fallback copies inside mc:AlternateContent; they should mirror
        # their Choice counterpart's ID.
        if _is_in_fallback(cNvPr):
            continue
        try:
            ids.append((int(cNvPr.get("id", "0")), cNvPr))
        except ValueError:
            ids.append((0, cNvPr))
    ids.sort(key=lambda x: x[0])
    for new_id, (_, cNvPr) in enumerate(ids, start=1):
        cNvPr.set("id", str(new_id))
    # Sync each fallback cNvPr ID with its Choice counterpart.
    for ac in slide.element.iter(f"{{{MC_NS}}}AlternateContent"):
        choice_sp = ac.find(f"{{{MC_NS}}}Choice/{{{P_NS}}}sp")
        fallback_sp = ac.find(f"{{{MC_NS}}}Fallback/{{{P_NS}}}sp")
        if choice_sp is None or fallback_sp is None:
            # Also handle group shapes and connector shapes.
            choice_sp = ac.find(f"{{{MC_NS}}}Choice/*")
            fallback_sp = ac.find(f"{{{MC_NS}}}Fallback/*")
        if choice_sp is not None and fallback_sp is not None:
            choice_cNvPr = choice_sp.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
            fallback_cNvPr = fallback_sp.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
            if choice_cNvPr is not None and fallback_cNvPr is not None:
                fallback_cNvPr.set("id", choice_cNvPr.get("id"))


def _is_in_fallback(el):
    """Return True if el is inside an mc:Fallback element."""
    parent = el.getparent()
    while parent is not None:
        if parent.tag == f"{{{MC_NS}}}Fallback":
            return True
        parent = parent.getparent()
    return False


def reorder_presentation(prs: "Presentation") -> None:
    """Fix the one known ``<p:presentation>`` child-order issue without reshuffling.

    Some template decks place ``<p:embeddedFontLst>`` after ``<p:defaultTextStyle>``
    or ``<p:extLst>``, which both the OOXML validator and PowerPoint dislike.
    This helper moves ``<p:embeddedFontLst>`` to just before the first of those
    two elements if it is out of place, leaving the rest of the child order
    (including the source deck's ``sldIdLst``/``notesMasterIdLst`` arrangement)
    unchanged so PowerPoint does not treat the file as repaired.

    Args:
        prs: A python-pptx ``Presentation`` object.

    Returns:
        None.
    """
    el = prs.part._element
    children = list(el)
    tags = [c.tag.split("}")[-1] for c in children]

    try:
        emb_idx = tags.index("embeddedFontLst")
    except ValueError:
        return

    target_names = ("defaultTextStyle", "extLst")
    target_indices = [i for i, t in enumerate(tags) if t in target_names]
    if not target_indices:
        return

    first_target = min(target_indices)
    if emb_idx < first_target:
        return

    emb = children[emb_idx]
    el.remove(emb)
    el.insert(first_target, emb)


def remove_layout_chrome(prs: "Presentation") -> None:
    """Remove footer/slide-number/date/header shapes from layouts/masters.

    The generator emits these as normal slide shapes so they are selectable.
    Leaving them on the layout would create duplicate locked background copies.

    Args:
        prs: A python-pptx ``Presentation`` object.

    Returns:
        None. Layout and master ``spTree`` elements are modified in place.
    """
    SLIDE_NUMBER_MARKER = "\u2039#\u203a"
    EMU_PER_INCH = 914400

    def _sp_text(sp):
        txBody = sp.find(f"{{{P_NS}}}txBody")
        if txBody is None:
            return ""
        return "".join(t.text or "" for t in txBody.findall(f".//{{{A_NS}}}t"))

    def _placeholder_type(sp):
        nv = sp.find(f"{{{P_NS}}}nvSpPr")
        if nv is None:
            return None
        nvPr = nv.find(f"{{{P_NS}}}nvPr")
        if nvPr is None:
            return None
        ph = nvPr.find(f"{{{P_NS}}}ph")
        return ph.get("type") if ph is not None else None

    def _shape_y_inches(sp):
        spPr = sp.find(f"{{{P_NS}}}spPr")
        if spPr is None:
            return 0.0
        xfrm = spPr.find(f"{{{A_NS}}}xfrm")
        if xfrm is None:
            return 0.0
        off = xfrm.find(f"{{{A_NS}}}off")
        if off is None:
            return 0.0
        try:
            return int(off.get("y", 0)) / EMU_PER_INCH
        except (ValueError, TypeError):
            return 0.0

    # Detect footer text by finding a short bottom-of-slide string repeated across
    # non-placeholder shapes on layouts/masters.
    text_counts = {}
    for part in [l for m in prs.slide_masters for l in m.slide_layouts] + list(prs.slide_masters):
        spTree = part._element.find(f"{{{P_NS}}}cSld/{{{P_NS}}}spTree")
        if spTree is None:
            continue
        for sp in spTree.findall(f"{{{P_NS}}}sp"):
            if _placeholder_type(sp):
                continue
            if _shape_y_inches(sp) < 6.0:
                continue
            text = _sp_text(sp).strip()
            if not text or text == SLIDE_NUMBER_MARKER:
                continue
            text_counts[text] = text_counts.get(text, 0) + 1

    footer_text = None
    candidates = [(t, c) for t, c in text_counts.items() if c >= 2 and len(t) <= 120]
    if candidates:
        footer_text = sorted(candidates, key=lambda x: (-x[1], len(x[0])))[0][0]

    for part in [l for m in prs.slide_masters for l in m.slide_layouts] + list(prs.slide_masters):
        spTree = part._element.find(f"{{{P_NS}}}cSld/{{{P_NS}}}spTree")
        if spTree is None:
            continue
        for sp in list(spTree.findall(f"{{{P_NS}}}sp")):
            ph_type = _placeholder_type(sp)
            if ph_type in ("ftr", "sldNum", "dt", "hdr"):
                spTree.remove(sp)
                continue
            text = _sp_text(sp).strip()
            if text == SLIDE_NUMBER_MARKER:
                if _shape_y_inches(sp) >= 6.0:
                    spTree.remove(sp)
                continue
            if footer_text and text == footer_text and _shape_y_inches(sp) >= 6.0:
                spTree.remove(sp)


# Alias kept for backwards compatibility with any project files that still call it.
remove_layout_footers_and_slide_numbers = remove_layout_chrome


def _has_math_xml(paras):
    """Return True if any paragraph run contains native Office Math XML."""
    for para in paras:
        if isinstance(para, str):
            continue
        if not isinstance(para, dict):
            continue
        runs = para.get("runs") if isinstance(para.get("runs"), list) else (
            para.get("text") if isinstance(para.get("text"), list) else []
        )
        for r in runs:
            if isinstance(r, dict) and r.get("math_xml"):
                return True
    return False


def _wrap_math_shape(shape):
    """Wrap a shape containing Office Math in mc:AlternateContent for compatibility.

    The fallback is a full clone of the original shape with the Office Math blocks
    stripped out. A minimal fallback shape caused PowerPoint to treat the wrapper
    as corrupt and remove math shapes, especially when they were nested inside
    group shapes.
    """
    sp = shape._element
    parent = sp.getparent()
    if parent is None:
        return

    index = parent.index(sp)

    nsmap = {"mc": MC_NS, "a14": A14_NS}
    ac = etree.Element(f"{{{MC_NS}}}AlternateContent", nsmap=nsmap)
    choice = etree.SubElement(ac, f"{{{MC_NS}}}Choice")
    choice.set("Requires", "a14")
    choice.append(sp)

    fallback = etree.SubElement(ac, f"{{{MC_NS}}}Fallback")
    fb_sp = deepcopy(sp)

    # Remove Office Math from the fallback so it is a plain, valid shape.
    for math_el in list(fb_sp.iter(f"{{{A14_NS}}}m")):
        math_parent = math_el.getparent()
        if math_parent is not None:
            math_parent.remove(math_el)

    # The fallback shape represents the same object as the Choice shape, so it
    # should carry the same cNvPr id/name. renumber_slide_ids will keep them in
    # sync and skip the fallback copy.
    fb_cNvPr = fb_sp.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}")
    if fb_cNvPr is not None:
        fb_cNvPr.set("name", fb_cNvPr.get("name", "MathShape"))

    fallback.append(fb_sp)
    parent.insert(index, ac)


# Generic shape helper
# ---------------------------------------------------------------------------

def _shapes(container):
    """Return the shapes collection for a slide or group shape."""
    return container.shapes


def _apply_shape_name(shape, name):
    """Set the non-visible shape name (cNvPr/@name)."""
    if name:
        try:
            shape.name = name
        except Exception:
            pass


def _apply_slide_number(shape, n):
    """Replace a shape's text with the dynamic slide number, preserving run formatting."""
    _apply_shape_name(shape, "SlideNumber")
    text = str(n)
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = text
        return
    p = tf.paragraphs[0]
    runs = list(p.runs)
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        p.text = text


def add_shape(
    slide_or_group: "Slide | GroupShape",
    kind: MSO_SHAPE | str | None,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    text: str | list = '',
    paragraphs: list | None = None,
    fill: str | dict | bool | None = None,
    line: str | bool | None = None,
    line_width: float | None = None,
    line_dash: str = "solid",
    line_head: dict | str | None = None,
    line_tail: dict | str | None = None,
    line_cap: str = "flat",
    line_cmpd: str = "sng",
    effects: list[dict] | None = None,
    rotation: float = 0,
    flip_h: bool = False,
    flip_v: bool = False,
    style: dict | None = None,
    adjustments: list[tuple[str, str]] | None = None,
    font: str | None = None,
    size: int | float | None = None,
    color: str | dict | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    underline: bool | str | None = None,
    strike: bool | str | None = None,
    highlight: str | dict | None = None,
    baseline: int | None = None,
    hyperlink: str | None = None,
    align: PP_ALIGN | str | None = None,
    anchor: MSO_ANCHOR | str | None = None,
    vertical_anchor: MSO_ANCHOR | str | None = None,
    margins: dict | None = None,
    wrap: bool | None = None,
    autofit: str | None = None,
    line_spacing: int | str | None = None,
    space_before: int | float | None = None,
    space_after: int | float | None = None,
    bullet: str | dict | None = None,
    pitch_family: str | None = None,
    charset: str | None = None,
    name: str | None = None,
    slide_number: int | None = None,
) -> "Shape":
    """Add a preset shape to a slide or group.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        kind: Preset shape name (e.g. ``'oval'``, ``'roundRect'``, ``'chevron'``,
            ``'rightArrow'``). Use ``None`` for a plain rectangle.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        text: Plain text string, or a list of paragraph strings/dicts. Ignored if
            ``paragraphs`` is supplied.
        paragraphs: List of paragraph dicts that fully describe the text frame.
            If provided, overrides ``text``.
        fill: Fill color. May be a hex string (``'FF0000'``), a theme token
            (``'theme_accent1'``, ``'scheme:accent1'``), ``'none'``/``False``,
            or a dict describing a gradient/pattern fill.
        line: Line color. Hex string (``'FF0000'``), theme token
            (``'theme_accent1'``, ``'scheme:accent1'``), or ``'none'``/``False``.
            ``None`` leaves python-pptx's default line. Only solid colors are
            supported for lines.
        line_width: Line width in points.
        line_dash: Dash style, e.g. ``'solid'``, ``'dash'``, ``'dot'``.
        line_head: Arrowhead specification. Either a string such as
            ``'arrow'``/``'none'`` or a dict with keys ``type``, ``w``, ``len``.
        line_tail: Same as ``line_head`` for the tail end.
        line_cap: Line cap style, e.g. ``'flat'``, ``'rnd'``, ``'sq'``.
        line_cmpd: Compound line style, e.g. ``'sng'``, ``'dbl'``.
        effects: List of effect dicts such as ``{"type": "outerShdw", ...}``.
        rotation: Rotation in degrees.
        flip_h: Flip horizontally.
        flip_v: Flip vertically.
        style: Theme style reference dict with keys like ``lnRef``,
            ``fillRef``, ``effectRef``, ``fontRef``.
        adjustments: List of ``(name, formula)`` tuples for preset-geometry
            adjustments, e.g. ``[("adj", "val 16667")]``.
        font: Font family name.
        size: Font size in points.
        color: Default text color (hex/theme/preset dict).
        bold: Default bold flag.
        italic: Default italic flag.
        underline: Default underline. ``True``/``False`` or a string such as
            ``'sng'``, ``'dbl'``.
        strike: Default strikethrough. ``True``/``False`` or a string such as
            ``'sngStrike'``.
        highlight: Text highlight color (hex or dict with alpha).
        baseline: Baseline offset in hundred-thousandths of an em (OpenXML
            ``a:baseline`` units).
        hyperlink: URL string applied to the default run.
        align: Paragraph alignment, e.g. ``'PP_ALIGN.LEFT'`` or ``PP_ALIGN.LEFT``.
        anchor: Vertical text anchor, e.g. ``'MSO_ANCHOR.MIDDLE'``.
        vertical_anchor: Alias for ``anchor``. If provided, takes precedence.
        margins: Dict with keys ``'lIns'``, ``'rIns'``, ``'tIns'``, ``'bIns'``
            in inches.
        wrap: Whether text wraps. ``None`` defaults to the python-pptx default.
        autofit: Autofit mode, e.g. ``'spAutoFit'``, ``'noAutofit'``,
            ``'normAutofit'``.
        line_spacing: Line spacing in thousandths of a line
            (``100000`` = single spacing). May also be a string ending in
            ``'pts'`` for point spacing.
        space_before: Paragraph spacing before in points.
        space_after: Paragraph spacing after in points.
        bullet: ``'none'``, ``'char'``, ``'autoNum'``, or a dict with ``type``
            and ``char``/``style``.
        pitch_family: Font pitch family string.
        charset: Font charset string.
        name: Semantic non-visible shape name (e.g. ``'Footer'``, ``'Date'``).
        slide_number: If set, replaces the shape text with this integer and
            marks the shape as the slide number.

    Returns:
        The created ``Shape`` object.
    """
    shape_kind = _shape_kind(kind)
    shape = _shapes(slide_or_group).add_shape(shape_kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _suppress_default_shadow(shape)
    _apply_style(shape, style)
    # Open/line shapes such as braces should not carry a solid fill; doing so
    # can trigger PowerPoint repair on some builds.
    if kind in ("rightBrace", "leftBrace"):
        fill = "none"
    _apply_fill(shape, fill)
    _apply_line(shape, line, line_width, line_dash, line_head, line_tail, line_cap, line_cmpd)
    _apply_effects(shape, effects)
    _apply_adjustments(shape, adjustments)
    _apply_rotation_and_flip(shape, rotation, flip_h, flip_v)

    # Preserve the plain "line" preset used by many hand-authored decks; python-pptx
    # defaults to "lineInv" for its LINE_INVERSE enum.
    if kind == "line":
        prstGeom = shape._element.spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", "line")

    paras = paragraphs if paragraphs is not None else (text if isinstance(text, list) else [text])
    if paras and not (len(paras) == 1 and paras[0] == ""):
        _apply_text_frame(
            shape.text_frame, paras,
            font=font, size=size, color=color, bold=bold, italic=italic,
            underline=underline, strike=strike, highlight=highlight, baseline=baseline,
            hyperlink=hyperlink, align=align, anchor=anchor, vertical_anchor=vertical_anchor,
            margins=margins, wrap=wrap, autofit=autofit, line_spacing=line_spacing,
            space_before=space_before, space_after=space_after, bullet=bullet,
            pitch_family=pitch_family, charset=charset,
        )
        if _has_math_xml(paras):
            _wrap_math_shape(shape)
    if slide_number is not None:
        _apply_slide_number(shape, slide_number)
    elif name is not None:
        _apply_shape_name(shape, name)
    return shape


def add_box(
    slide_or_group: "Slide | GroupShape",
    text: str | list,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | dict | bool | None = d.COL["blue"],
    color: str | dict | None = None,
    line: str | bool | None = None,
    line_width: float | None = None,
    font: str | None = None,
    size: int | float | None = None,
    bold: bool | None = None,
    align: PP_ALIGN | str | None = None,
    rounded: bool = True,
    rounded_adj: int | float | None = None,
    anchor: MSO_ANCHOR | str | None = MSO_ANCHOR.MIDDLE,
    style: dict | None = None,
    pitch_family: str | None = None,
    charset: str | None = None,
    wrap: bool | None = None,
    name: str | None = None,
    slide_number: int | None = None,
    **kwargs,
) -> "Shape":
    """Add a colored box. ``text`` may be a string or a list of paragraphs.

    When the box has no fill and no line, this delegates to :func:`add_text` so
    PowerPoint treats it as a text box (``txBox="1"``) rather than a rectangle.
    This avoids repair/lock issues for title-like text on some builds.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        text: Text string or list of paragraph strings/dicts.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        fill: Fill color (hex/theme/gradient/pattern dict) or ``'none'``.
            Defaults to ``d.COL["blue"]``.
        color: Default text color (hex/theme/preset dict).
        line: Border color (hex string or theme token) or ``'none'``.
        line_width: Border width in points.
        font: Font family name.
        size: Font size in points.
        bold: Default bold flag.
        align: Paragraph alignment, e.g. ``'PP_ALIGN.CENTER'``.
        rounded: Whether to use a rounded rectangle.
        rounded_adj: Rounding adjustment value (preset-geometry ``adj`` formula
            numeric value). Larger values produce a more pronounced corner.
        anchor: Vertical text anchor, e.g. ``'MSO_ANCHOR.MIDDLE'``.
        style: Theme style reference dict.
        pitch_family: Font pitch family string.
        charset: Font charset string.
        wrap: Whether text wraps.
        name: Semantic shape name.
        slide_number: If set, displays this integer as the slide number.
        **kwargs: Extra text/formatting arguments passed to :func:`add_text` or
            :func:`add_shape`, such as ``rotation``, ``effects``, ``margins``,
            ``line_spacing``, etc.

    Returns:
        The created ``Shape`` (or text-box shape when fill/line are absent).
    """
    if fill in (None, False, "none") and line in (None, False, "none"):
        # Shape-only kwargs (rotation, flip, effects) are not text-frame
        # properties; apply them to the text box after creating it.
        shape_kwargs = {}
        for k in ("rotation", "flip_h", "flip_v", "effects"):
            if k in kwargs:
                shape_kwargs[k] = kwargs.pop(k)
        box = add_text(
            slide_or_group, text, x, y, w, h,
            font=font, size=size, color=color, bold=bold, align=align,
            anchor=anchor, pitch_family=pitch_family, charset=charset,
            wrap=wrap, name=name, slide_number=slide_number,
            **kwargs,
        )
        if shape_kwargs:
            _apply_rotation_and_flip(
                box,
                shape_kwargs.get("rotation", 0),
                shape_kwargs.get("flip_h", False),
                shape_kwargs.get("flip_v", False),
            )
            if "effects" in shape_kwargs:
                _apply_effects(box, shape_kwargs["effects"])
        return box
    kind = "roundRect" if rounded else "rect"
    adjustments = None
    if rounded and rounded_adj is not None:
        adjustments = [("adj", f"val {rounded_adj}")]
    return add_shape(
        slide_or_group, kind, x, y, w, h,
        text=text,
        fill=fill, line=line, line_width=line_width,
        font=font, size=size, color=color, bold=bold, align=align,
        anchor=anchor, style=style, pitch_family=pitch_family, charset=charset,
        wrap=wrap, adjustments=adjustments,
        name=name, slide_number=slide_number,
        **kwargs,
    )


def add_text(
    slide_or_group: "Slide | GroupShape",
    text: str | list,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font: str | None = None,
    size: int | float | None = None,
    color: str | dict | None = None,
    bold: bool | None = None,
    align: PP_ALIGN | str | None = None,
    fill: str | dict | bool | None = None,
    line: str | bool | None = None,
    line_width: float | None = None,
    anchor: MSO_ANCHOR | str | None = MSO_ANCHOR.MIDDLE,
    vertical_anchor: MSO_ANCHOR | str | None = None,
    pitch_family: str | None = None,
    charset: str | None = None,
    margins: dict | None = None,
    wrap: bool | None = None,
    autofit: str | None = None,
    line_spacing: int | str | None = None,
    space_before: int | float | None = None,
    space_after: int | float | None = None,
    bullet: str | dict | None = None,
    name: str | None = None,
    slide_number: int | None = None,
    **kwargs,
) -> "Shape":
    """Add a text box with multiple paragraphs and bullets.

    When ``fill`` and ``line`` are both ``None`` and the target is a slide (not
    a group), a native text-box shape is created. Otherwise a rectangle shape
    with the supplied fill/line is used.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        text: Text string or list of paragraph strings/dicts.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        font: Font family name.
        size: Font size in points.
        color: Default text color (hex/theme/preset dict).
        bold: Default bold flag.
        align: Paragraph alignment, e.g. ``'PP_ALIGN.LEFT'``.
        fill: Fill color for the backing rectangle, if any.
        line: Border color for the backing rectangle, if any.
        line_width: Border width in points (only used with a backing rectangle).
        anchor: Vertical text anchor, e.g. ``'MSO_ANCHOR.MIDDLE'``.
        vertical_anchor: Alias for ``anchor``. If provided, takes precedence.
        pitch_family: Font pitch family string.
        charset: Font charset string.
        margins: Dict with keys ``'lIns'``, ``'rIns'``, ``'tIns'``, ``'bIns'``
            in inches.
        wrap: Whether text wraps. ``None`` defaults to ``True`` for native text
            boxes and leaves rectangle shapes untouched.
        autofit: Autofit mode, e.g. ``'spAutoFit'`` or ``'noAutofit'``.
        line_spacing: Line spacing in thousandths (``100000`` = single) or a
            string ending in ``'pts'`` for point spacing.
        space_before: Paragraph spacing before in points.
        space_after: Paragraph spacing after in points.
        bullet: ``'none'``, ``'char'``, ``'autoNum'``, or a dict.
        name: Semantic shape name.
        slide_number: If set, displays this integer as the slide number.
        **kwargs: Additional arguments forwarded to :func:`add_shape` (e.g.
            ``rotation``, ``effects``, ``italic``, ``underline``).

    Returns:
        The created text-box or rectangle ``Shape``.
    """
    paragraphs = text if isinstance(text, list) else [text]
    shapes = _shapes(slide_or_group)
    is_group = not hasattr(shapes, 'add_textbox')
    if fill is None and line is None and not is_group:
        box = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        # python-pptx textboxes default to wrap="none"; most source text boxes
        # use square wrap (no explicit bodyPr/@wrap), so default to True.
        effective_wrap = True if wrap is None else wrap
        _apply_text_frame(
            box.text_frame, paragraphs,
            font=font, size=size, color=color, bold=bold, align=align,
            anchor=anchor if vertical_anchor is None else None,
            vertical_anchor=vertical_anchor,
            pitch_family=pitch_family, charset=charset,
            margins=margins, wrap=effective_wrap, autofit=autofit, line_spacing=line_spacing,
            space_before=space_before, space_after=space_after, bullet=bullet,
            **kwargs,
        )
        if _has_math_xml(paragraphs):
            _wrap_math_shape(box)
        if slide_number is not None:
            _apply_slide_number(box, slide_number)
        elif name is not None:
            _apply_shape_name(box, name)
        return box

    return add_shape(
        slide_or_group, "rect", x, y, w, h,
        paragraphs=paragraphs,
        fill=fill, line=line, line_width=line_width,
        font=font, size=size, color=color, bold=bold, align=align,
        anchor=anchor, vertical_anchor=vertical_anchor,
        pitch_family=pitch_family, charset=charset,
        margins=margins, wrap=wrap, autofit=autofit, line_spacing=line_spacing,
        space_before=space_before, space_after=space_after, bullet=bullet,
        name=name, slide_number=slide_number,
        **kwargs,
    )


def add_label(
    slide_or_group: "Slide | GroupShape",
    text: str | list,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font: str | None = None,
    size: int | float | None = None,
    color: str | dict | None = None,
    bold: bool | None = None,
    align: PP_ALIGN | str | None = None,
    anchor: MSO_ANCHOR | str | None = MSO_ANCHOR.MIDDLE,
    pitch_family: str | None = None,
    charset: str | None = None,
    name: str | None = None,
    slide_number: int | None = None,
    **kwargs,
) -> "Shape":
    """Add a text-only label.

    This is a convenience wrapper around :func:`add_text` that passes
    ``fill=None`` and ``line=None``. On slides this typically creates a
    transparent native text box; on group shapes it falls back to the
    rectangle-path behavior of :func:`add_text`.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        text: Text string or list of paragraph strings/dicts.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        font: Font family name.
        size: Font size in points.
        color: Text color (hex/theme/preset dict).
        bold: Bold flag.
        align: Paragraph alignment.
        anchor: Vertical text anchor.
        pitch_family: Font pitch family string.
        charset: Font charset string.
        name: Semantic shape name.
        slide_number: If set, displays this integer as the slide number.
        **kwargs: Extra arguments forwarded to :func:`add_text`.

    Returns:
        The created text-box ``Shape``.
    """
    return add_text(
        slide_or_group, text, x, y, w, h,
        font=font, size=size, color=color, bold=bold, align=align,
        anchor=anchor, pitch_family=pitch_family, charset=charset,
        name=name, slide_number=slide_number,
        **kwargs,
    )


# ---------------------------------------------------------------------------
# Lines and arrows
# ---------------------------------------------------------------------------

def add_line(
    slide_or_group: "Slide | GroupShape",
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str | dict | None = d.COL["sep"],
    width: float = 0.75,
    dash: str = "solid",
    head: dict | str | None = None,
    tail: dict | str | None = None,
    cap: str = "flat",
    cmpd: str = "sng",
    style: dict | None = None,
    rotation: float = 0,
) -> "Shape":
    """Draw a straight line shape from ``(x1, y1)`` to ``(x2, y2)``.

    ``(x1, y1)`` .. ``(x2, y2)`` describe the line's *unrotated* extent and
    direction; ``rotation`` is then applied around its center (OOXML semantics).

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        x1: Start X position in inches.
        y1: Start Y position in inches.
        x2: End X position in inches.
        y2: End Y position in inches.
        color: Line color (hex/theme/preset dict).
        width: Line width in points.
        dash: Dash style, e.g. ``'solid'`` or ``'dash'``.
        head: Arrowhead at the line end. String such as ``'arrow'`` or a dict
            with ``type``, ``w``, ``len``.
        tail: Arrowhead at the line start.
        cap: Line cap style, e.g. ``'flat'``, ``'rnd'``, ``'sq'``.
        cmpd: Compound line style, e.g. ``'sng'`` or ``'dbl'``.
        style: Theme style reference dict.
        rotation: Rotation in degrees, applied around the line's center.

    Returns:
        The created line ``Shape``.
    """
    x, y = min(x1, x2), min(y1, y2)
    w, h = abs(x2 - x1), abs(y2 - y1)
    shape = add_shape(
        slide_or_group, "line", x, y, w, h,
        fill=None, line=color, line_width=width, line_dash=dash,
        line_head=head, line_tail=tail, line_cap=cap, line_cmpd=cmpd,
        style=style, rotation=rotation,
    )
    # Orient the line to point from (x1,y1) to (x2,y2).
    spPr = shape._element.spPr
    xfrm = spPr.xfrm
    if x2 < x1:
        xfrm.set("flipH", "1")
    if y2 < y1:
        xfrm.set("flipV", "1")
    return shape


def add_arrow(
    slide_or_group: "Slide | GroupShape",
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str | dict | None = d.COL["blue"],
    width: float = 2,
    dash: str = "solid",
    head: dict | str | None = "arrow",
    tail: dict | str | None = None,
) -> "Shape":
    """Draw a straight arrow with optional head/tail.

    This is a convenience wrapper around :func:`add_line` that defaults the
    head end to an arrow.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        x1: Start X position in inches.
        y1: Start Y position in inches.
        x2: End X position in inches.
        y2: End Y position in inches.
        color: Arrow line color (hex/theme/preset dict).
        width: Line width in points.
        dash: Dash style.
        head: Arrowhead at the end (default ``'arrow'``). Set to ``None`` or
            ``'none'`` to remove.
        tail: Arrowhead at the start.

    Returns:
        The created arrow ``Shape``.
    """
    return add_line(
        slide_or_group, x1, y1, x2, y2,
        color=color, width=width, dash=dash, head=head, tail=tail,
    )


def add_connector(
    slide_or_group: "Slide | GroupShape",
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str | dict | None = d.COL["sep"],
    width: float = 0.75,
    dash: str = "solid",
    head: dict | str | None = None,
    tail: dict | str | None = None,
    cap: str = "flat",
    cmpd: str = "sng",
    style: dict | None = None,
    kind: str = "straight",
    rotation: float = 0,
    preset: str | None = None,
    adjustments: list | None = None,
) -> "Connector":
    """Add a connector shape (``cxnSp``).

    ``(x1, y1)`` .. ``(x2, y2)`` describe the connector's *unrotated* bounding
    box and direction (the endpoint ordering encodes flipH/flipV). ``rotation``
    is then applied around the box center, matching OOXML semantics, so a
    rotated elbow/bent connector lands in the right place and orientation.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        x1: Start X position in inches.
        y1: Start Y position in inches.
        x2: End X position in inches.
        y2: End Y position in inches.
        color: Connector line color (hex/theme/preset dict).
        width: Line width in points.
        dash: Dash style.
        head: Arrowhead at the end.
        tail: Arrowhead at the start.
        cap: Line cap style.
        cmpd: Compound line style.
        style: Theme style reference dict.
        kind: Connector geometry: ``'straight'``, ``'elbow'``, or ``'curved'``.
        rotation: Rotation in degrees, applied around the box center.
        preset: Exact preset geometry name (e.g. ``'bentConnector2'``) to
            reproduce faithfully. Overrides the geometry ``kind`` maps to.
        adjustments: Optional list of ``(name, fmla)`` guide tuples for the
            preset's ``avLst``.

    Returns:
        The created connector shape.
    """
    kind_map = {
        "straight": MSO_CONNECTOR.STRAIGHT,
        "elbow": MSO_CONNECTOR.ELBOW,
        "curved": MSO_CONNECTOR.CURVE,
    }
    conn_kind = kind_map.get(kind, MSO_CONNECTOR.STRAIGHT)
    conn = _shapes(slide_or_group).add_connector(
        conn_kind,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    _apply_style(conn, style)
    _apply_line(conn, color, width, dash, head, tail, cap, cmpd)
    if preset:
        _set_connector_preset(conn, preset, adjustments)
    elif adjustments:
        # No preset override, but the source carried adjustment guides (e.g. a
        # bent connector's ``adj1`` bend position) that must be applied to the
        # default preset — otherwise the turning point snaps back to center.
        _apply_adjustments(conn, adjustments)
    if rotation:
        conn.rotation = rotation
    return conn


def connect_shapes(conn, begin=None, end=None) -> None:
    """Attach a connector's ends to shapes via ``stCxn``/``endCxn``.

    Unlike python-pptx's ``begin_connect``/``end_connect``, this does *not*
    reposition the connector — it keeps the connector's already-correct baked
    geometry and only records the connection sites, so a viewer that reroutes
    connected connectors (PowerPoint, LibreOffice) reproduces the original
    bend exactly.

    Args:
        conn: A python-pptx ``Connector``.
        begin: ``(shape, idx)`` for the start, or ``None``.
        end: ``(shape, idx)`` for the end, or ``None``.
    """
    cxnPr = conn._element.nvCxnSpPr.cNvCxnSpPr
    # cNvCxnSpPr children must be ordered <a:stCxn> then <a:endCxn>.
    for tag, spec in (("a:endCxn", end), ("a:stCxn", begin)):
        for existing in cxnPr.findall(qn(tag)):
            cxnPr.remove(existing)
        if spec is None:
            continue
        shape, idx = spec
        el = cxnPr.makeelement(qn(tag), {"id": str(shape.shape_id), "idx": str(idx)})
        cxnPr.insert(0, el)


def _set_connector_preset(conn, preset: str, adjustments: list | None = None) -> None:
    """Override a connector's preset geometry (and adjustment guides).

    python-pptx only exposes STRAIGHT/ELBOW/CURVE, which map to a fixed set of
    presets (e.g. ELBOW -> ``bentConnector3``). Source decks use finer variants
    such as ``bentConnector2``; this rewrites ``prstGeom`` to match.
    """
    prstGeom = conn._element.spPr.find(qn("a:prstGeom"))
    if prstGeom is None:
        return
    prstGeom.set("prst", preset)
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        avLst = prstGeom.makeelement(qn("a:avLst"), {})
        prstGeom.append(avLst)
    for gd in list(avLst):
        avLst.remove(gd)
    for name, fmla in (adjustments or []):
        avLst.append(prstGeom.makeelement(qn("a:gd"), {"name": name, "fmla": fmla}))


def add_callout(
    slide: "Slide",
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    kind: str = "blue",
    size: int | float = 16,
) -> "Shape":
    """Add a pre-styled callout box.

    Args:
        slide: A python-pptx ``Slide``.
        text: Callout text.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        kind: Callout style: ``'blue'``, ``'green'``, ``'red'``, or ``'yellow'``.
        size: Font size in points.

    Returns:
        The created callout ``Shape``.
    """
    styles = {
        "blue": (d.COL["blue"], d.COL["white"], d.COL["blue"]),
        "green": (d.COL["green"], d.COL["black"], d.COL["green"]),
        "red": (d.COL["dark_red"], d.COL["white"], d.COL["black"]),
        "yellow": (d.COL["yellow"], d.COL["black"], d.COL["black"]),
    }
    fill, color, line = styles.get(kind, styles["blue"])
    return add_box(slide, text, x, y, w, h, fill=fill, color=color, line=line,
                   line_width=2.25 if kind != "yellow" else 6, size=size)


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------

def _svg_to_png(svg_path: Path, width_px: int, height_px: int) -> Path:
    """Rasterize an SVG to a temporary PNG using rsvg-convert.

    Falls back to cairosvg's SVG->PNG conversion if available, then Inkscape.
    Raises RuntimeError if no converter is available.
    """
    png_path = Path(tempfile.mktemp(suffix=".png"))
    converters = [
        ("rsvg-convert", [
            "rsvg-convert", "-w", str(width_px), "-h", str(height_px),
            "-o", str(png_path), str(svg_path)
        ]),
        ("cairosvg", None),
        ("inkscape", [
            "inkscape", str(svg_path),
            "--export-filename=" + str(png_path),
            "--export-width=" + str(width_px),
            "--export-height=" + str(height_px),
        ]),
    ]
    for name, cmd in converters:
        if name == "cairosvg":
            try:
                import cairosvg
                cairosvg.svg2png(
                    url=str(svg_path),
                    write_to=str(png_path),
                    output_width=width_px,
                    output_height=height_px,
                )
                return png_path
            except Exception:
                continue
        if cmd and shutil.which(name):
            try:
                subprocess.run(cmd, check=True, capture_output=True)
                if png_path.exists() and png_path.stat().st_size > 0:
                    return png_path
            except Exception:
                continue
    raise RuntimeError(
        f"Cannot rasterize SVG {svg_path}: install rsvg-convert (librsvg), "
        "cairosvg, or inkscape."
    )


def _convert_to_png(
    source: Path,
    width_px: int | None = None,
    height_px: int | None = None,
) -> Path:
    """Convert an image asset to a temporary PNG for python-pptx ingestion.

    Handles SVG (vector rasterization), WDP/JPEG-XR/HD Photo (ImageMagick),
    and WebP/HEIC/HEIF or any other format Pillow can decode. SVG is rendered
    at the requested pixel size; raster formats keep their original pixel
    dimensions and let python-pptx scale them to the slide coordinates.

    Args:
        source: Path to the image asset.
        width_px: Target width in pixels for SVG rasterization. Ignored for
            raster formats.
        height_px: Target height in pixels for SVG rasterization. Ignored for
            raster formats.

    Returns:
        Path to a temporary PNG file.

    Raises:
        RuntimeError: If no converter is available or conversion fails.
    """
    suffix = source.suffix.lower()
    png_path = Path(tempfile.mktemp(suffix=".png"))

    if suffix == ".svg":
        return _svg_to_png(source, width_px or 0, height_px or 0)

    if suffix == ".wdp":
        convert_cmd = shutil.which("convert")
        if convert_cmd:
            try:
                subprocess.run(
                    [convert_cmd, str(source), str(png_path)],
                    check=True,
                    capture_output=True,
                )
                if png_path.exists() and png_path.stat().st_size > 0:
                    return png_path
            except Exception:
                pass
        raise RuntimeError(
            f"Cannot convert WDP image {source}: install ImageMagick "
            "(`convert`) or pre-convert the asset to PNG/JPEG."
        )

    # WebP, HEIC, HEIF, and any other format Pillow can decode.
    try:
        # pillow-heif registers itself as a Pillow plugin on import.
        try:
            import pillow_heif  # noqa: F401
        except Exception:
            pass
        with PILImage.open(source) as im:
            # Preserve transparency where possible; otherwise flatten to RGB.
            if im.mode in ("RGBA", "P", "LA"):
                im = im.convert("RGBA")
            else:
                im = im.convert("RGB")
            im.save(png_path, "PNG")
        if png_path.exists() and png_path.stat().st_size > 0:
            return png_path
    except Exception as e:
        raise RuntimeError(
            f"Cannot convert image {source} to PNG: {e}. "
            "Install a suitable decoder (e.g. pillow-heif for HEIC/HEIF) "
            "or convert the asset to PNG/JPEG before building."
        ) from e

    raise RuntimeError(f"Cannot convert image {source} to PNG: unknown format.")


def add_image(
    slide_or_group: "Slide | GroupShape",
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    crop: dict | None = None,
    lum: dict | None = None,
    rotation: float = 0,
) -> "Picture":
    """Add an image from the assets directory.

    Image files must live under the configured assets directory
    (:func:`set_assets_dir`). Native formats (PNG, JPEG, GIF, BMP, TIFF, WMF,
    EMF) are passed straight through. SVG is rasterized to PNG at 300 dpi.
    WebP, HEIC/HEIF, and WDP are converted to PNG on the fly when the needed
    decoder/tool is installed.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        name: Asset filename (e.g. ``'logo.png'`` or ``'diagram.svg'``).
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        crop: Optional crop dict with keys ``'l'``, ``'r'``, ``'t'``, ``'b'``.
            Values are hundred-thousandths of the image dimension (e.g.
            ``100000`` crops the full edge).
        lum: Reserved/non-functional. The implementation attempts to set
            ``pic.brightness`` and ``pic.contrast``, but these attributes are
            not available on python-pptx ``Picture`` objects in this
            environment, so the values are silently ignored.
        rotation: Clockwise rotation in degrees, applied about the picture's
            box center (matching OOXML ``xfrm`` rotation). Used to reproduce
            rotated freeform/custom-geometry artwork whose flip is baked into
            the SVG content, so directional shapes (arrows) keep their heading.

    Returns:
        The created ``Picture`` shape.

    Raises:
        FileNotFoundError: If the asset name is empty or the file does not exist.
        RuntimeError: If a non-native asset cannot be converted to PNG.
    """
    if not name:
        raise FileNotFoundError(
            "Image asset name is empty. "
            "The generator could not resolve this picture to an asset file; "
            "replace with a TODO or place the image in assets/."
        )
    path = ASSETS / name
    if not path.exists():
        raise FileNotFoundError(
            f"Image asset not found: {path}. "
            f"Place the image in the assets/ directory or correct the filename."
        )

    suffix = path.suffix.lower()
    if suffix == ".svg":
        # Render at 300 dpi for crisp output; the picture is then scaled to the
        # requested w/h by add_picture().
        dpi = 300
        picture_path = _convert_to_png(path, int(w * dpi), int(h * dpi))
        _tmp_png = picture_path
    elif suffix not in _NATIVE_IMAGE_EXTS:
        # Convert once at original resolution; add_picture() scales to w/h.
        picture_path = _convert_to_png(path)
        _tmp_png = picture_path
    else:
        picture_path = path
        _tmp_png = None

    try:
        pic = _shapes(slide_or_group).add_picture(
            str(picture_path), Inches(x), Inches(y), Inches(w), Inches(h)
        )
    finally:
        if _tmp_png is not None:
            try:
                _tmp_png.unlink()
            except Exception:
                pass

    # Preserve a description on the picture so round-trips keep a semantic
    # link to the asset name instead of falling back to generic filenames.
    try:
        cNvPr = pic._element.find(f"{{{P_NS}}}nvPicPr/{{{P_NS}}}cNvPr")
        if cNvPr is not None:
            cNvPr.set("descr", Path(name).stem)
    except Exception:
        pass

    if rotation:
        pic.rotation = rotation

    if crop:
        try:
            if crop.get("l") is not None:
                pic.crop_left = int(crop["l"]) / 100000
            if crop.get("r") is not None:
                pic.crop_right = int(crop["r"]) / 100000
            if crop.get("t") is not None:
                pic.crop_top = int(crop["t"]) / 100000
            if crop.get("b") is not None:
                pic.crop_bottom = int(crop["b"]) / 100000
        except (ValueError, TypeError, AttributeError):
            pass
    if lum:
        try:
            if lum.get("bright") is not None:
                pic.brightness = int(lum["bright"]) / 1000
            if lum.get("contrast") is not None:
                pic.contrast = int(lum["contrast"]) / 1000
        except Exception:
            pass
    return pic


CHART_TYPE_MAP = {
    "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
    "COLUMN_PERCENT_STACKED": XL_CHART_TYPE.COLUMN_STACKED_100,
    "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
    "BAR_STACKED": XL_CHART_TYPE.BAR_STACKED,
    "BAR_PERCENT_STACKED": XL_CHART_TYPE.BAR_STACKED_100,
    "LINE_CLUSTERED": XL_CHART_TYPE.LINE,
    "LINE_STACKED": XL_CHART_TYPE.LINE_STACKED,
    "LINE_PERCENT_STACKED": XL_CHART_TYPE.LINE_STACKED_100,
    "LINE_MARKERS": XL_CHART_TYPE.LINE_MARKERS,
    "PIE": XL_CHART_TYPE.PIE,
    "PIE_EXPLODED": XL_CHART_TYPE.PIE_EXPLODED,
    "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
    "AREA_CLUSTERED": XL_CHART_TYPE.AREA,
    "AREA_STACKED": XL_CHART_TYPE.AREA_STACKED,
    "AREA_PERCENT_STACKED": XL_CHART_TYPE.AREA_STACKED_100,
}


def add_chart(
    slide: "Slide",
    chart_type: str,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list,
    series: list[dict],
    title: str | None = None,
) -> "GraphicFrame":
    """Add a chart using python-pptx ``ChartData``.

    Args:
        slide: A python-pptx ``Slide``.
        chart_type: Chart type key, e.g. ``'COLUMN_CLUSTERED'``, ``'BAR_STACKED'``,
            ``'LINE_MARKERS'``, ``'PIE'``. Unknown keys fall back to a clustered
            column chart.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        categories: List of category labels.
        series: List of series dicts, each with ``'name'`` and ``'values'``.
        title: Optional chart title.

    Returns:
        The ``GraphicFrame`` containing the chart.
    """
    chart_data = ChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s["name"], s["values"])
    xl_type = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    graphic_frame = slide.shapes.add_chart(
        xl_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    )
    chart = graphic_frame.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    return graphic_frame


def add_movie(
    slide: "Slide",
    video_name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    poster_name: str | None = None,
) -> "Movie":
    """Add a video/movie shape with an optional poster-frame image.

    Video and poster files must live under the configured assets directory
    (:func:`set_assets_dir`). If ``poster_name`` is provided but missing, it is
    silently ignored.

    Args:
        slide: A python-pptx ``Slide``.
        video_name: Asset filename of the video.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        poster_name: Optional asset filename for the poster frame image.

    Returns:
        The created movie shape.

    Raises:
        FileNotFoundError: If the video asset is not found.
    """
    video_path = ASSETS / video_name
    if not video_path.exists():
        raise FileNotFoundError(f"Video asset not found: {video_path}")
    poster_path = str(ASSETS / poster_name) if poster_name else None
    if poster_path and not Path(poster_path).exists():
        poster_path = None
    movie = slide.shapes.add_movie(
        str(video_path),
        Inches(x), Inches(y), Inches(w), Inches(h),
        poster_frame_image=poster_path,
    )
    return movie


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------

def _make_border_ln(edge_tag, color, width_pt, dash):
    """Build an <a:lnL/R/T/B> element representing a table cell border."""
    ln = etree.Element(qn(edge_tag))
    ln.set("w", str(int(Pt(width_pt))))
    if _is_theme(color):
        sf = etree.SubElement(ln, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:schemeClr"))
        sc.set("val", color.replace("theme_", ""))
    else:
        sf = etree.SubElement(ln, qn("a:solidFill"))
        c = etree.SubElement(sf, qn("a:srgbClr"))
        c.set("val", color.lstrip("#").upper())
    dash_style = dash if dash else "solid"
    pd = etree.SubElement(ln, qn("a:prstDash"))
    pd.set("val", dash_style)
    for attr, val in (("headEnd", "none"), ("tailEnd", "none")):
        he = etree.SubElement(ln, qn(f"a:{attr}"))
        he.set("type", val)
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    return ln


_TCPR_CHILD_ORDER = [
    "lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr",
    "noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill",
    "effectLst", "effectDag", "scene3d", "sp3d", "extLst",
]

_TCPR_BORDER_TAGS = {qn(f"a:{t}") for t in _TCPR_CHILD_ORDER[:6]}


def _reorder_tcPr(tcPr):
    """Sort <a:tcPr> children into schema order (borders, then fills/effects)."""
    order = {tag: idx for idx, tag in enumerate(_TCPR_CHILD_ORDER)}

    def sort_key(child):
        tag = child.tag.split("}")[-1]
        return order.get(tag, len(order))

    tcPr[:] = sorted(tcPr, key=sort_key)


def set_cell_border(
    cell: "Cell",
    edge: str,
    color: str,
    width_pt: float = 0.5,
    dash: str = "solid",
) -> None:
    """Set a single border edge of a table cell.

    Borders must precede fills inside ``<a:tcPr>``; this helper inserts the new
    border and then reorders the children to satisfy the schema.

    Args:
        cell: A python-pptx table ``Cell``.
        edge: Edge to style: ``'left'``, ``'right'``, ``'top'``, or ``'bottom'``.
        color: Border color as a hex string (``'FF0000'``) or theme token
            (``'theme_accent1'``, ``'scheme:accent1'``). Only solid colors are
            supported for cell borders.
        width_pt: Border width in points.
        dash: Dash style, e.g. ``'solid'`` or ``'dash'``.

    Returns:
        None.

    Raises:
        ValueError: If ``edge`` is not one of the supported edge names.
    """
    edge_tag = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}.get(edge)
    if edge_tag is None:
        raise ValueError(f"Invalid edge: {edge}")
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn(edge_tag)):
        tcPr.remove(old)
    ln = _make_border_ln(edge_tag, color, width_pt, dash)
    tcPr.append(ln)
    _reorder_tcPr(tcPr)


def set_table_shape_size(table: "Table", w_inches: float, h_inches: float) -> None:
    """Override the outer graphic-frame size of a table.

    Args:
        table: A python-pptx ``Table`` object.
        w_inches: New width in inches.
        h_inches: New height in inches.

    Returns:
        None.
    """
    table._graphic_frame._element.xfrm.ext.cx = Inches(w_inches)
    table._graphic_frame._element.xfrm.ext.cy = Inches(h_inches)


def _apply_cell_text(cell, text, font=None, size=None, color=None, bold=None,
                     italic=None, underline=None, strike=None, highlight=None,
                     baseline=None, align=None, anchor=None, pitch_family=None, charset=None):
    """Apply text and formatting to a table cell."""
    if isinstance(text, list):
        paragraphs = text
    else:
        paragraphs = [{"text": text}]
    tf = cell.text_frame
    _apply_text_frame(
        tf, paragraphs,
        font=font, size=size, color=color, bold=bold, italic=italic,
        underline=underline, strike=strike, highlight=highlight, baseline=baseline,
        align=align, anchor=anchor, vertical_anchor=anchor,
        pitch_family=pitch_family, charset=charset,
    )


def add_custom_table(
    slide_or_group: "Slide | GroupShape",
    rows: list[list] | None = None,
    fills: list[list] | None = None,
    x: float = 0,
    y: float = 0,
    w: float = 1,
    h: float = 1,
    *,
    size: int | float = 14,
    header_fill: str | dict | None = None,
    header_color: str | dict | None = None,
    body_color: str | dict | None = None,
    border_color: str | None = None,
    row_heights: list[float] | None = None,
    col_widths: list[float] | None = None,
    borders: bool = True,
    cells: list[list[dict]] | None = None,
    border_width: float = 0.5,
    border_dash: str = "solid",
) -> "Table":
    """Add a table with explicit per-cell control.

    Two call styles are supported:

    1. Legacy::

           rows=[["A", "B"], ["C", "D"]]
           fills=[["FFFFFF", "CCCCCC"], ["EEEEEE", "DDDDDD"]]

    2. Full cell specs::

           cells=[[{"text": "A", "fill": "FFFFFF"}, ...], ...]

    In the full style, each cell dict may include ``text``, ``fill``,
    ``font``, ``size``, ``color``, ``bold``, ``italic``, ``underline``,
    ``strike``, ``highlight``, ``baseline``, ``align``, ``anchor``,
    ``margins`` (dict with ``marL``/``marR``/``marT``/``marB`` in inches),
    ``borders`` (dict mapping edge to ``{"color", "w", "dash"}``),
    ``border_color``, ``border_width``, ``border_dash``, ``gridSpan`` and
    ``rowSpan`` for merges.

    Args:
        slide_or_group: A python-pptx ``Slide`` or ``GroupShape``.
        rows: Legacy 2-D list of cell text.
        fills: Legacy 2-D list of fill colors (must align with ``rows``).
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        size: Default font size in points.
        header_fill: Reserved for future use; currently not applied automatically.
        header_color: Reserved for future use; currently not applied automatically.
        body_color: Reserved for future use; currently not applied automatically.
        border_color: Default border color for cells.
        row_heights: Optional list of row heights in inches.
        col_widths: Optional list of column widths in inches.
        borders: If ``True``, draw all four borders on every cell using the
            default border color/width/dash.
        cells: 2-D list of per-cell spec dicts. If provided, ``rows`` and
            ``fills`` are ignored.
        border_width: Default border width in points.
        border_dash: Default border dash style.

    Returns:
        The created python-pptx ``Table``.

    Raises:
        ValueError: If neither ``rows`` nor ``cells`` is provided.
    """
    if cells is not None:
        n_rows = len(cells)
        n_cols = max(len(r) for r in cells)
    elif rows is not None:
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        cells = None
    else:
        raise ValueError("Either rows or cells must be provided")

    table = _shapes(slide_or_group).add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h)).table

    if row_heights:
        for ri, rh in enumerate(row_heights):
            if ri < len(table.rows):
                table.rows[ri].height = Inches(rh)
    if col_widths:
        for ci, cw in enumerate(col_widths):
            if ci < len(table.columns):
                table.columns[ci].width = Inches(cw)

    for ri in range(n_rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            if cells is not None:
                spec = cells[ri][ci] if ci < len(cells[ri]) else {}
            else:
                spec = {"text": rows[ri][ci] if ci < len(rows[ri]) else ""}
                if fills and ri < len(fills) and ci < len(fills[ri]):
                    spec["fill"] = fills[ri][ci]

            text = spec.get("text", "")
            _apply_cell_text(
                cell, text,
                font=spec.get("font"),
                size=spec.get("size"),
                color=spec.get("color"),
                bold=spec.get("bold"),
                italic=spec.get("italic"),
                underline=spec.get("underline"),
                strike=spec.get("strike"),
                highlight=spec.get("highlight"),
                baseline=spec.get("baseline"),
                align=spec.get("align"),
                anchor=spec.get("anchor"),
                pitch_family=spec.get("pitch_family"),
                charset=spec.get("charset"),
            )

            fill = spec.get("fill")
            if fill is not None:
                _apply_fill(cell, fill)

            anchor = spec.get("anchor")
            if anchor is not None:
                if isinstance(anchor, str) and anchor.startswith("MSO_ANCHOR."):
                    anchor = getattr(MSO_VERTICAL_ANCHOR, anchor.split(".")[-1], anchor)
                cell.vertical_anchor = anchor

            margins = spec.get("margins")
            if margins:
                for key, attr in (("marL", "left"), ("marR", "right"), ("marT", "top"), ("marB", "bottom")):
                    val = margins.get(key)
                    if val is not None:
                        try:
                            setattr(cell, f"margin_{attr}", Inches(float(val)))
                        except Exception:
                            pass

            cell_borders = spec.get("borders")
            if cell_borders:
                bc = spec.get("border_color", border_color or d.COL["sep"])
                bw = spec.get("border_width", border_width)
                bd = spec.get("border_dash", border_dash)
                for edge, b in cell_borders.items():
                    set_cell_border(cell, edge, b.get("color", bc), b.get("w", bw), b.get("dash", bd))
            elif borders:
                bc = spec.get("border_color", border_color or d.COL["sep"])
                bw = spec.get("border_width", border_width)
                bd = spec.get("border_dash", border_dash)
                for edge in ("left", "right", "top", "bottom"):
                    set_cell_border(cell, edge, bc, bw, bd)

    # Apply merges after all cells are populated.
    if cells is not None:
        for ri in range(n_rows):
            for ci in range(n_cols):
                spec = cells[ri][ci] if ci < len(cells[ri]) else {}
                grid_span = int(spec.get("gridSpan") or 1)
                row_span = int(spec.get("rowSpan") or 1)
                if grid_span > 1 or row_span > 1:
                    end_ri = ri + row_span - 1
                    end_ci = ci + grid_span - 1
                    if end_ri < n_rows and end_ci < n_cols:
                        table.cell(ri, ci).merge(table.cell(end_ri, end_ci))

    return table


# ---------------------------------------------------------------------------
# Groups
# ---------------------------------------------------------------------------

def add_group(slide: "Slide", x: float, y: float, w: float, h: float) -> "GroupShape":
    """Add a group shape and return it so children can be added to it.

    The ``x``, ``y``, ``w`` and ``h`` arguments are accepted for API symmetry
    but are currently ignored by the implementation. Use :func:`set_group_bounds`
    after adding children to set the group's bounds.

    Args:
        slide: A python-pptx ``Slide``.
        x: Reserved left position in inches (ignored).
        y: Reserved top position in inches (ignored).
        w: Reserved width in inches (ignored).
        h: Reserved height in inches (ignored).

    Returns:
        The created ``GroupShape``.
    """
    group = slide.shapes.add_group_shape()
    return group


def set_group_bounds(
    group: "GroupShape",
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    """Set the outer position and size of a group after children are added.

    Args:
        group: A python-pptx ``GroupShape``.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.

    Returns:
        None.
    """
    xfrm = group._element.grpSpPr.xfrm
    for attr in ("rot", "flipH", "flipV"):
        xfrm.attrib.pop(attr, None)
    xfrm.find(qn("a:off")).set("x", str(Emu(Inches(x))))
    xfrm.find(qn("a:off")).set("y", str(Emu(Inches(y))))
    xfrm.find(qn("a:ext")).set("cx", str(Emu(Inches(w))))
    xfrm.find(qn("a:ext")).set("cy", str(Emu(Inches(h))))
    xfrm.find(qn("a:chOff")).set("x", "0")
    xfrm.find(qn("a:chOff")).set("y", "0")
    xfrm.find(qn("a:chExt")).set("cx", str(Emu(Inches(w))))
    xfrm.find(qn("a:chExt")).set("cy", str(Emu(Inches(h))))


# ---------------------------------------------------------------------------
# Background and notes
# ---------------------------------------------------------------------------

def add_background(slide: "Slide", fill: str | dict) -> None:
    """Set the slide background from a fill spec.

    Args:
        slide: A python-pptx ``Slide``.
        fill: Fill spec. May be a hex color string (``'FF0000'``), a theme
            color token (``'theme_accent1'``), or a dict describing a gradient
            or solid fill::

                {"type": "gradient", "angle": 90,
                 "stops": [(0, "FFFFFF"), (1, "000000")]}
                {"type": "solid", "color": "FF0000"}

    Returns:
        None.
    """
    if isinstance(fill, dict):
        if fill.get("type") == "gradient":
            stops = fill.get("stops", [])
            if len(stops) < 2:
                stops = [(0, "FFFFFF"), (1, "000000")]
            slide.background.fill.gradient()
            for i, (pos, color) in enumerate(stops):
                if i >= len(slide.background.fill.gradient_stops):
                    break
                gs = slide.background.fill.gradient_stops[i]
                gs.position = pos
                _apply_theme_color(gs.color, color)
        elif fill.get("type") == "solid":
            slide.background.fill.solid()
            _apply_theme_color(slide.background.fill.fore_color, fill["color"])
    else:
        slide.background.fill.solid()
        _apply_theme_color(slide.background.fill.fore_color, fill)


def add_notes(slide: "Slide", text: str) -> None:
    """Replace the notes text for a slide.

    Args:
        slide: A python-pptx ``Slide``.
        text: Notes text.

    Returns:
        None.
    """
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def set_slide_hidden(slide: "Slide", hidden: bool = True) -> None:
    """Hide or unhide a slide in the presentation.

    Slide visibility is the ``show`` attribute on the slide's ``<p:sld>``
    element: ``show="0"`` marks it hidden (skipped in slideshows and greyed
    out in the slide sorter). python-pptx exposes no property for this, so the
    attribute is set directly on ``slide.element``. Visible is the absence of
    the attribute, so unhiding removes it rather than writing ``show="1"``.

    Args:
        slide: A python-pptx ``Slide``.
        hidden: ``True`` to hide the slide, ``False`` to unhide it.

    Returns:
        None.
    """
    if hidden:
        slide.element.set("show", "0")
    elif slide.element.get("show") is not None:
        del slide.element.attrib["show"]


def postprocess_svg_fallbacks(pptx_path: Path | str) -> None:
    """Remove raster fallbacks from SVG pictures in slide layouts/masters.

    PowerPoint on some builds repairs decks that contain a picture with both a
    raster ``<a:blip>`` and an ``<asvg:svgBlip>`` extension. Leaving only the
    SVG blip and adding a Default content type for ``.svg`` prevents that
    repair dialog.

    Args:
        pptx_path: Path to the generated ``.pptx`` file.

    Returns:
        None. The file is rewritten in place.
    """
    pptx_path = Path(pptx_path)
    tmp_path = pptx_path.with_suffix(".pptx.tmp")

    A_BLIP = qn("a:blip")
    R_EMBED = qn("r:embed")
    ASVG_SVGBLIP = "{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip"

    with zipfile.ZipFile(pptx_path, "r") as zin:
        items = {name: zin.read(name) for name in zin.namelist()}

    # Ensure .svg has a Default content type instead of per-part Overrides.
    ct_name = "[Content_Types].xml"
    if ct_name in items and any(name.endswith(".svg") for name in items):
        ct = items[ct_name].decode("utf-8")
        if '<Default Extension="svg"' not in ct:
            ct = ct.replace(
                '<Default Extension="xml"',
                '<Default Extension="svg" ContentType="image/svg+xml"/><Default Extension="xml"',
            )
        # Remove any part-specific SVG overrides now covered by the default.
        ct = re.sub(
            r'<Override PartName="/ppt/media/[^"]+\.svg" ContentType="image/svg\+xml"/>',
            "",
            ct,
        )
        items[ct_name] = ct.encode("utf-8")

    layout_xml_names = {
        name for name in items
        if name.startswith("ppt/slideLayouts/slideLayout") and name.endswith(".xml")
    } | {
        name for name in items
        if name.startswith("ppt/slideMasters/slideMaster") and name.endswith(".xml")
    }

    rels_to_remove = {}  # rels_path -> set of rIds

    for xml_name in layout_xml_names:
        root = etree.fromstring(items[xml_name])
        changed = False
        for blip in root.iter(A_BLIP):
            svg_blip = blip.find(".//" + ASVG_SVGBLIP)
            if svg_blip is None:
                continue
            main_embed = blip.get(R_EMBED)
            if not main_embed:
                continue
            # The main blip is the raster fallback; drop it and keep the SVG.
            del blip.attrib[R_EMBED]
            changed = True
            rels_name = os.path.dirname(xml_name) + "/_rels/" + os.path.basename(xml_name) + ".rels"
            if rels_name in items:
                rels_to_remove.setdefault(rels_name, set()).add(main_embed)
        if changed:
            items[xml_name] = etree.tostring(
                root, encoding="UTF-8", standalone=True, xml_declaration=True
            )

    for rels_name, drop_ids in rels_to_remove.items():
        rels_root = etree.fromstring(items[rels_name])
        for rel in list(rels_root):
            if rel.get("Id") in drop_ids:
                # Only drop if it points to a raster image, not the SVG.
                target = rel.get("Target", "").lower()
                if target.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                    rels_root.remove(rel)
        items[rels_name] = etree.tostring(
            rels_root, encoding="UTF-8", standalone=True, xml_declaration=True
        )

    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in items.items():
            zout.writestr(name, data)

    tmp_path.replace(pptx_path)


def postprocess_powerpoint_native(pptx_path: Path | str) -> None:
    """Normalize serialization details that trigger PowerPoint's repair dialog.

    python-pptx emits single-quoted XML declarations and a zip entry order that
    differs from what PowerPoint writes. Rewriting declarations to double quotes
    and placing slide parts near the start of the package avoids the
    "PowerPoint found a problem with content" repair prompt on some builds.

    Args:
        pptx_path: Path to the generated ``.pptx`` file.

    Returns:
        None. The file is rewritten in place.
    """
    pptx_path = Path(pptx_path)
    tmp_path = pptx_path.with_suffix(".pptx.native.tmp")

    def _zip_priority(name: str) -> tuple:
        if name == "[Content_Types].xml":
            return (0, name)
        if name == "_rels/.rels":
            return (1, name)
        if name.startswith("ppt/slides/"):
            return (2, name)
        if name.startswith("ppt/notesSlides/"):
            return (3, name)
        if name.startswith("ppt/slideLayouts/"):
            return (4, name)
        if name.startswith("ppt/slideMasters/"):
            return (5, name)
        if name.startswith("ppt/notesMasters/"):
            return (6, name)
        if name.startswith("ppt/theme/"):
            return (7, name)
        if name.startswith("ppt/media/"):
            return (8, name)
        return (9, name)

    with zipfile.ZipFile(pptx_path, "r") as zin:
        items = [(name, zin.read(name)) for name in zin.namelist()]

    processed = []
    for name, data in items:
        if name.endswith(".xml") or name.endswith(".rels"):
            text = data.decode("utf-8")
            # Single-quoted XML declaration -> double-quoted.
            if text.startswith("<?xml version='"):
                text = text.replace("<?xml version='1.0' encoding='UTF-8' standalone='yes'?>",
                                    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>')
            data = text.encode("utf-8")
        processed.append((name, data))

    processed.sort(key=lambda x: _zip_priority(x[0]))

    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in processed:
            zout.writestr(name, data)

    tmp_path.replace(pptx_path)



