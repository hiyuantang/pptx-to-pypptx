#!/usr/bin/env python3
"""Generate (overwrite) selected slide .py files from a target PPTX.

This is the code-generation half of the round-trip workflow. It always fully
overwrites the requested slide files; it never does partial edits.

Usage:
    uv run python generate_slides.py \
        --target "path/to/target.pptx" \
        --project-dir my-deck \
        --slides 4,5,9

--slides is required and accepts a single number, a range (2-5), or a comma
list (3,7,9). There is no --all option.
"""

import argparse
import re
import sys
import tempfile
import zipfile
from pathlib import Path

# Make scaffold importable when running this script directly.
_SCRIPT_DIR = Path(__file__).resolve().parent
if str(_SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPT_DIR))

from pptx import Presentation

from helpers.assets import load_media_map, sync_assets
from helpers.slide_codegen import (
    generate_slide_code,
    generate_layout_chrome_code,
    detect_footer_text,
)
from helpers.slide_xml import load_slide_layout_path
from helpers.slides import get_slide_title, sanitize_name


def parse_slide_arg(arg: str, total: int) -> list[int]:
    """Parse '1,3,5' or '2-4' into a sorted list of 1-based slide numbers."""
    result = set()
    for part in arg.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            try:
                result.update(range(int(start), int(end) + 1))
            except ValueError:
                raise ValueError(f"Invalid slide range: {part!r}")
        else:
            try:
                result.add(int(part))
            except ValueError:
                raise ValueError(f"Invalid slide number: {part!r}")
    return sorted(n for n in result if 1 <= n <= total)


def count_slides(pptx: Path) -> int:
    with zipfile.ZipFile(pptx, "r") as zf:
        return len([
            n for n in zf.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        ])


def find_existing_slide_file(project_dir: Path, idx: int) -> Path | None:
    """Find the slide file for a given 1-based slide index."""
    slides_dir = project_dir / "slides"
    candidates = sorted(slides_dir.glob(f"s{idx:02d}_*.py"))
    return candidates[0] if candidates else None


def generate_slides(target: Path, project_dir: Path, slides: list[int]):
    if not target.exists():
        raise FileNotFoundError(f"Target PPTX not found: {target}")
    project_dir = Path(project_dir)
    slides_dir = project_dir / "slides"
    if not slides_dir.exists():
        raise FileNotFoundError(f"Slides directory not found: {slides_dir}")

    total = count_slides(target)
    slide_numbers = [s for s in slides if 1 <= s <= total]
    if not slide_numbers:
        print("No valid slide numbers to generate.")
        return

    # Sync new/changed media into the project's assets/ folder, deduplicated
    # by content hash, then load the mapping from raw media name -> asset name.
    sync_assets(target, project_dir)
    media_names = load_media_map(project_dir)

    # Open the target presentation once so we can read each slide's actual
    # layout index from python-pptx's layout collection.
    prs = Presentation(str(target))
    # Flatten layouts across ALL slide masters. ``prs.slide_layouts`` only
    # exposes the first master's layouts, so decks with multiple masters would
    # fail the lookup for every slide bound to a non-first master and fall back
    # to index 0 (often the Title Slide layout, dragging its branding onto every
    # slide). build_deck.py resolves LAYOUT against the same flattened order.
    layouts = [layout for master in prs.slide_masters for layout in master.slide_layouts]

    def _layout_idx(slide_idx: int) -> int:
        layout = prs.slides[slide_idx - 1].slide_layout
        for i, candidate in enumerate(layouts):
            if candidate == layout:
                return i
        return 0

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        with zipfile.ZipFile(target, "r") as zf:
            zf.extractall(tmp_path)

        slide_dir = tmp_path / "ppt" / "slides"

        def _slide_key(path: Path):
            m = re.search(r"(\d+)", path.stem)
            return int(m.group(1)) if m else 0

        slide_files = sorted(slide_dir.glob("slide*.xml"), key=_slide_key)

        layout_paths = [
            load_slide_layout_path(p) for p in slide_files
        ]
        unique_layouts = [p for p in set(layout_paths) if p]
        footer_text = detect_footer_text(unique_layouts) if unique_layouts else None

        for idx in slide_numbers:
            slide_xml = slide_files[idx - 1]
            title = get_slide_title(slide_xml)
            stem = sanitize_name(title) or f"slide_{idx}"
            new_name = f"s{idx:02d}_{stem}.py"
            new_path = slides_dir / new_name

            existing = find_existing_slide_file(project_dir, idx)
            if existing and existing != new_path:
                existing.unlink()

            layout_idx = _layout_idx(idx)
            assets_dir = project_dir / "assets"
            body = generate_slide_code(slide_xml, media_names, title, assets_dir=assets_dir)

            layout_xml = layout_paths[idx - 1]
            if layout_xml:
                chrome = generate_layout_chrome_code(
                    layout_xml, media_names, idx, footer_text, assets_dir=assets_dir
                )
                if chrome:
                    body = f"{body}\n{chrome}"

            new_path.write_text(
                f'from lib import shapes\n'
                f'from pptx.enum.text import PP_ALIGN, MSO_ANCHOR\n'
                f'from pptx.dml.color import RGBColor\n\n'
                f'TITLE = {title!r}\n'
                f'LAYOUT = {layout_idx}\n\n'
                f'def add_slide(prs, slide, n):\n'
                f'{body}\n',
                encoding="utf-8",
            )
            print(f"Generated {new_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate slide .py files from a target PPTX"
    )
    parser.add_argument("--target", required=True, help="Target PPTX file")
    parser.add_argument(
        "--project-dir", required=True, help="Project directory containing slides/"
    )
    parser.add_argument(
        "--slides",
        required=True,
        help="Slide numbers to generate: 4 | 2-5 | 3,7,9",
    )
    args = parser.parse_args()

    total = count_slides(Path(args.target))
    slides = parse_slide_arg(args.slides, total)
    generate_slides(Path(args.target), Path(args.project_dir), slides)


if __name__ == "__main__":
    main()
