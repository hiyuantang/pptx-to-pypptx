"""Shared .pptx utilities: slide counting and slide-range parsing.

Internal helper imported by the scripts in the parent directory; not run directly.
"""

import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET


P = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_PER_INCH = 914400


def count_slides(pptx: Path) -> int:
    """Count slide*.xml entries inside a .pptx zip."""
    with zipfile.ZipFile(pptx, "r") as zf:
        return sum(
            1 for n in zf.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )


def read_slide_size(pptx: Path) -> tuple[float, float]:
    """Return the presentation canvas width and height in inches."""
    with zipfile.ZipFile(pptx, "r") as zf:
        try:
            presentation_xml = zf.read("ppt/presentation.xml")
        except KeyError as exc:
            raise ValueError("PPTX is missing ppt/presentation.xml") from exc
    root = ET.fromstring(presentation_xml)
    slide_size = root.find(f"{{{P}}}sldSz")
    if slide_size is None:
        raise ValueError("PPTX presentation has no slide-size declaration")
    try:
        width = int(slide_size.get("cx")) / EMU_PER_INCH
        height = int(slide_size.get("cy")) / EMU_PER_INCH
    except (TypeError, ValueError) as exc:
        raise ValueError("PPTX slide-size declaration is invalid") from exc
    if width <= 0 or height <= 0:
        raise ValueError("PPTX slide size must be positive")
    return width, height


def write_base_deck(source: Path, dest: Path) -> None:
    """Write a *base* deck: the template shell with all content slides removed.

    Keeps every slide master, layout, and the theme so ``build_deck.py`` can bind
    generated slides to the right layout, but drops the source deck's own slides.
    The result is a minimal, self-contained shell the project builds from, so a
    plain ``build_deck.py`` needs no source ``.pptx`` at build time.
    """
    from pptx import Presentation

    prs = Presentation(str(source))
    for slide in list(prs.slides):
        slide_id = slide.slide_id
        rId = None
        for rel in prs.part.rels.values():
            if rel.target_part == slide.part:
                rId = rel.rId
                break
        if rId:
            prs.part.drop_rel(rId)
        for sldId in list(prs.slides._sldIdLst):
            if sldId.id == slide_id:
                prs.slides._sldIdLst.remove(sldId)
                break
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))


def parse_slide_range(arg: str, total: int, *, allow_all: bool = True) -> list[int]:
    """Parse a slide selection into a sorted list of 1-based slide numbers.

    Accepts a single number (``14``), a range (``8-12``), a comma list
    (``4,5,9``), or ``all`` (only when ``allow_all`` is set). Numbers are clamped
    to ``[1, total]``. Raises ``ValueError`` on malformed input, or on ``all``
    when ``allow_all`` is False.
    """
    if arg.strip().lower() == "all":
        if not allow_all:
            raise ValueError(
                "'all' is not supported here; specify slides like 4 | 2-5 | 3,7,9"
            )
        return list(range(1, total + 1))

    result = set()
    for part in arg.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start, end = part.split("-", 1)
            try:
                result.update(range(int(start), int(end) + 1))
            except ValueError:
                raise ValueError(f"Invalid slide range: {part!r}")
        else:
            try:
                result.add(int(part))
            except ValueError:
                raise ValueError(f"Invalid slide number: {part!r}")
    return sorted(n for n in result if 1 <= n <= total)
