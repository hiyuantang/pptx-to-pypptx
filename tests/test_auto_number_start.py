import sys
import types
import unittest
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))
sys.path.insert(0, str(ROOT / "template"))
design_stub = types.ModuleType("lib.design")
design_stub.COL = {"blue": "0000FF", "sep": "CCCCCC"}
sys.modules.setdefault("lib.design", design_stub)

from helpers.slide_codegen import _paragraph_props
from helpers.slide_model import normalize_paragraph
from helpers.slide_xml import A, _parse_para_level_props
from lib.shapes import _set_paragraph_bullet


class AutoNumberStartTests(unittest.TestCase):
    def test_auto_number_start_is_preserved_end_to_end(self):
        paragraph = ET.fromstring(
            f'''<a:pPr xmlns:a="{A}">
              <a:buAutoNum type="arabicPeriod" startAt="3"/>
            </a:pPr>'''
        )

        parsed = _parse_para_level_props(paragraph)
        normalized = normalize_paragraph({"text": "Third step", "runs": [], **parsed})
        generated = _paragraph_props(normalized)

        self.assertEqual(parsed["bullet_start_at"], 3)
        self.assertEqual(normalized["bullet_start_at"], 3)
        self.assertEqual(
            generated["bullet"],
            {"type": "autoNum", "style": "arabicPeriod", "start_at": 3},
        )

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        target = textbox.text_frame.paragraphs[0]
        _set_paragraph_bullet(target, generated["bullet"])

        auto_number = target._p.get_or_add_pPr().find(qn("a:buAutoNum"))
        self.assertIsNotNone(auto_number)
        self.assertEqual(auto_number.get("type"), "arabicPeriod")
        self.assertEqual(auto_number.get("startAt"), "3")


if __name__ == "__main__":
    unittest.main()
